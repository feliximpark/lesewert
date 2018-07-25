# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""
#%% Import Libraries
# IMPORT LIBRARIES

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib
import seaborn as sns
import datetime as dt

# Deutsch als mathematischen Standard einsetzen
import locale
locale.setlocale(locale.LC_ALL, "deu_deu")


# import dt. Mathezeichen (Komma statt punkt bei Dezimalzahlen)
matplotlib.rcParams['axes.formatter.use_locale'] = True


#import der Schrift für Matplotlib
from matplotlib import rcParams
import sys
import os
import matplotlib.font_manager as fm
#Schrift laden
fpath1 = 'C:\\Windows\\Fonts\\Campton-Light.otf'
campton_light = fm.FontProperties(fname=fpath1)


import timeit
import time
import requests
#für die Bildbearbeitung:
from PIL import Image
from PIL import ImageDraw, ImageFont
from io import StringIO
import textwrap

# alles von pptx importieren
from pptx import Presentation 
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Mm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import time 

# Glätten der Line-Plots
from scipy.interpolate import interp1d


#%% LISTEN UND ONJEKTE SETZEN

#%% Listen und Objekte


platzierung = ["AA", "ZA", "SK", "ÜA", "TS1", "TS2", "TS3", "TS4", "TS5"]
platzierung_dict = {"AA": "Aufmacher", "ZA": "Zweitaufmacher", "SK": "Seitenkeller", "ÜA": "Überaufmacher", "TS1":"Einspalter", 
                    "TS2":"Zweispalter", "TS3":"Dreispalter", "TS4":"Vierspalter", "TS5": "Fünfspalter"}
darstellung_dict = {"NA": "Nachricht", "AK": "Ankündigung", "SA": "Struktur. Artikel", "ZU": "Zitattext/Umfrage",\
                        "SE":"Serie", "GA": "Gastartikel", "FF": "Freies Format", 'LB':'Leserbrief','ES': 'Essay', \
                        'RK': 'Rezension/Kritik', 'KM': 'Komm./Kolumne', 'BB': 'Bildbericht', 'BN': 'Bildnachricht', \
                        'GT': 'Geschichtstext', 'KI': 'Karikatur/Illustration/Grafik', 'ST': 'Serviceelement/Tabelle', \
                        'HG': 'Hintergrundkasten', 'IV': 'Interview', "RP": "Reportage/Porträt", "BF":"Bericht/Feature", "BI": "Bildbe./Bildn." }
darstellungsform = ["BF", "RP", "IV", "KM", "BF", "LB", "SA", "NA", "RK", "ZU"]

seitentitel_dict = {'Veranstaltungen':'Veranstaltungen', 'Kirchen':"Kirchen", 
                    'Lokales':"Lokales", 'Biberach':"Biberach", 'Rund um Biberach':'Rund um Biberach', 
                    'Von der Schussen zur Rot': "Schussen zur Rot", 'Von der Rottum zur Iller':"Rottum zur Iller",
                    "Kultur": "Kultur", "Friedrichshafen":"Friedrichshafen", "Region": "Region", "Biberach / Service":"Biberach/Service", 
                    "Rund um Biberach / Schussen - Rot":"Biberach/Sch.-Rot",'Schussen-Rot / Rottum-Iller':"Sch.-Rot/Rottum-Ill.", 
                    'Wir in Kreis und Region':'Wir in Kreis + Region', 'Friedrichshafen / Service':'Friedrichsh./Service', 
                    'Markdorf und der westliche Bodenseekreis':"Markdorf/westl. BSK", 'Friedrichshafen / Service':'Friedrich./Service', 
                    'LokalesImmenstaad / Oberteuringen':'Immens./Obert.', 'Wir am See':"Wir am See", 
                    'Ravensburg / Weingarten':'Ravensburg/Weing.', 'Ravensburg':"Ravensburg", "Weingarten":"Weingarten", 
                    'Oberschwaben & Allgäu': 'Oberschw. & Allgäu', 'Langenargen / Eriskirch / Kressbronn':'Lang./Eris./Kressbr.', 
                    'Immenstaad / Oberteuringen / Meckenbeuren': 'Immenst./Obert./Meck.', 'Die Seite mit der Maus':"Seite mit der Maus", 
                    "Extraseiten":"Extraseiten", 'Rottum-Iller / Kreis & Region': r'Rottum-Iller/K+R',
                    'Langenargen / Eriskirch':"Langenargen/Eris.","Oberteuringen /Immenstaad/Meckenbeuren":"Ober./Imm./Mecken.", 
                     'Friedrichshafen / Tettnang':"Fried./Tettnang","Garten": "Garten", "Kreismusikfest":"Kreismusikfest",   
                     'Aktion Engagieren & kassieren':'Aktion Engagieren & kassieren', "Wissen": "Wissen", 
                     'Biberach / Rund um Biberach':'Rund um Biberach', "Panorama":"Panorama",
                     'Immenstaad / Meckenbeuren':'Imm. / Meckenb.', "Ernährung":"Ernährung", 
                     'Meine Heimat. Mein Verein. Extraseiten':"Heimat/Vereine", "Bauen & Wohnen": "Bauen & Wohnen", 
                     'Langenargen/Kressbronn':'Langenargen/Kress.', "Gemeinden":"Gemeinden",
                     "Umland":"Umland", "Lokalsport":"Lokalsport", 'Kultur Lokal/Service':'Kultur Lokal/Service', 
                     'rund um Biberach':'Rund um Biberach', 'Kultur Lokal':'Kultur Lokal', 'Immenstaad':'Immenstaad', 
                     'Reise & Erholung':'Reise & Erholung', "Auto & Verkehr":"Auto & Verkehr", 'Oberteuringen / Immenstaad':'Oberteuringen / Immenstaad', 
                     "Vermischtes":"Vermischtes", "Medien":"Medien", "Multimedia":"Multimedia", "Sternenhimmel":"Sternenhimmel",  
                     "Kino":"Kino", "Fernsehen & Freizeit":"Fernsehen/Freizeit", "Geld & Service":"Geld & Service",
                     "Fernsehen": "Fernsehen", "Leitartikel":"Leitartikel", "Familie":"Familie", "Von der Schussen zur Umlach":"Schussen z. Umlach", 
                     'Langenargen / Kressbronn':'Langenargen / Kressbronn', 
                     "Markdorf / Oberteuringen / Immenstaad / Salem":"Mark./Obert./Imm./Sal.", 
                     "Oberteuringen / Meckenbeuren":"Oberteuringen / Meckenbeuren", "Langenargen":"Langenargen", 
                     "Regionalsport / Lokalsport":"Reg./Lokalsport", "Regionalsport":"Regionalsport", "Szene":"Szene", 
                     "Mode":"Mode", "Literatur": "Literatur", "Tiere":"Tiere", "Fernsehen":"Fernsehen", 
                     "Lokale Eins BIB": "Lokale Eins BIB","Lokale Eins FHA": "Lokale Eins FHA", "Lokale Eins RV": "Lokale Eins RV", 
                     "Wochenende": "Wochenende", 'Menschen':"Menschen", 'Lebensart':"Lebensart", 'Unterhaltung':"Unterhaltung",
                     "Szene am Wochenende": "Szene am WE", 'Meine Seite':"Meine Seite" }
                    

ressort_list = ['Titel', 'Wir im Süden', 'Seite Drei', 'Nachrichten & Hintergrund',
       'Meinung & Dialog', 'Wirtschaft', 'Journal', 'Kultur',
       'Ratgeber', 'Panorama', 'Lokales', 'Lokalsport',
       'Sport', 'Wochenende']               

#mantel_ressorts = ['Titel', 'Wir im Süden', 'Seite Drei', 'Nachr. & Hint.',
       #'Mein. & Dialog', 'Wirtschaft', 'Journal', 'Kultur', 'Fernsehen',
       #'Reise', 'Ratgeber', 'Panorama',
       #'Vermischtes', 'Sport', 'Wochenende']     
mantel_ressorts = ['Titelseite', 'Politik', 'Seite 3',
       'Zwischen Weser und Rhein', 'Wirtschaft', 'Aus aller Welt', 'Sport']
     
ressort_dict={'Titel':'Titel', 'Wir im Süden':'Wir im Süden', 'Seite Drei':'Seite Drei', 'Nachrichten & Hintergrund':'Nachr. & Hint.',
       'Meinung & Dialog':'Mein. & Dialog' , 'Wirtschaft':'Wirtschaft', 'Journal':'Journal', 'Kultur':'Kultur',
       'Ratgeber':'Ratgeber', 'Panorama':'Panorama', 'Lokales':'Lokales', 'Lokalsport':'Lokalsport',
       'Sport':'Sport', 'Wochenende':'Wochenende', 'Fernsehen': "Fernsehen", 
       'Titelseite':"Titelseite", 'Tagesthema':"Tagesthema", 'Die Dritte Seite':"Die Dritte Seite", 'Politik':"Politik",
       'Aus aller Welt':"Aus aller Welt",'Wirtschaft':"Wirtschaft", 'Kultur': "Kultur", 'Sport':"Sport", 'Entdecken':"Entdecken",
       'Multimed. Reportage': "Multim. Reportage",'Leserforum':"Leserforum", 'Wochenende':"Wochenende", "Die Seite Drei":"Die Seite Drei", 
       "Panorama": "Panorama", "Wissenswert":"Wissenswert", "Unsere Leser und wir":"Unsere Leser + wir", 
       'Lokales Stuttgart STZ':"Lokales Stutt. STZ", 'Lokales/Region STN':"Lokales/Region STN",  
       'Region/Baden-Württ. STZ':"Region/BaWü STZ",  'Stuttgart 5. Buch':"Stuttgart 5. Buch",
                                    'Filderstadt 5. Buch':"Filderstadt 5. Buch",
                                    'Lokales Ludwigsburg':'Lokales Ludwigsburg', 'Lokalsport':"Lokalsport", 
                                     "Landesnachrichten":"Landesnachr."}                    

ausgaben_liste=[]     

seitentitel_lokal = []

liste_sportarten = ["Frauentennis", "Wasserball", 'American Football','Football', 
                    "Judo",'Leichtathletik',\
                   'Reitsport', 'Formel E', 'Fußball','Marathon', 'Badminton', 
                   'Golfsport','Snooker', 'Eishockey', \
                   'Langlauf', 'Radfahren', 'Handball', 'Skilanglauf', 'Tennis',
                   'Biathlon', 'Rudern','Rallye',\
                   'Laufen', 'Skispringen','Motorsport', 'Tischtennis','Radsport',
                   'Eiskunstlauf', 'Springsport', 'Ringen',\
                   'Basketball', 'Golf', 'Dressurreiten', 'Frauenfußball',
                   'Fünfkampf', 'Frauenvolleyball', 'Gewichtheben',\
                   'Rugby', 'Frauenhandball', 'Boxen', 'Volleyball', 'Fechten',
                   'Turnen','Beachvolleyball', "Kanu", 'Beachvolleball',
                   'Bouldern', 'Dressur', 'Formel 1', 'Kraftdreikampf', 
                   'Paralympics', 'Zehnkampf', 'Zeitfahren', 'Rudern']   

zeitung_attribute_dict={'Ihre Zeitung ist unabhängig':"unabhängig", 
               'Ihre Zeitung ist informativ': "informativ",
               'Ihre Zeitung ist niveauvoll':"niveauvoll", 
               'Ihre Zeitung ist interessant':"interessant",
               'Ihre Zeitung ist übersichtlich':"übersichtlich",
               'Ihre Zeitung ist verständlich geschrieben':"verständlich",
               'Ihre Zeitung ist zeitgemäß':"zeitgemäß",
               'Ihre Zeitung glaubwürdig':"glaubwürdig",
               'Ihre Zeitung ist glaubwürdig':"glaubwürdig", 
               'Ihre Zeitung ist nah am Leser':"nah am Leser", 
               'Ihre Zeitung ist unterhaltsam':"unterhaltsam"}


#UMFRAGE-LISTEN
# Liste um die Umfrageergebnisse "Welche Themen interessieren Sie..." zu ordnen
# und auf verschiedene Folien zu verteilen
umfrage_sport = ['Berichte über Fußball', 'Berichte über Basketball', 'Berichte über Biathlon',
       'Berichte über Handball','Berichte über Schwimmen oder Wasserspringen',
       'Berichte über Volleyball', 'Berichte über American Football',
       'Berichte über Bergsport', 'Berichte über Eishockey',
       'Berichte über Leichtathletik', 'Berichte über Motorsport',
       'Berichte über Pferdesport', 'Berichte über Radsport',
       'Berichte über Schach', 'Berichte über Tennis', 'Berichte über Boxen' ]

umfrage_sport_dict = {'Berichte über Fußball':"Fußball", 
                      'Berichte über Basketball':"Basketball", 
                      'Berichte über Biathlon':"Biathlon",
                      'Berichte über Handball':"Handball",
                      'Berichte über Schwimmen oder Wasserspringen':"Schwimmen",
                      'Berichte über Volleyball': 'Volleyball',
                      'Berichte über Bergsport':"Bergsport", 
                      'Berichte über Eishockey':"Eishockey",
                      'Berichte über Leichtathletik':"Leichtathletik",
                      'Berichte über Motorsport':"Motorsport",
                      'Berichte über Pferdesport':"Pferdesport", 
                      'Berichte über Radsport':"Radsport",
                      'Berichte über Schach':"Schach",
                      'Berichte über Tennis':"Tennis",
                      'Berichte über Boxen':"Boxen"}

umfrage_gesell_dict = {"Berichte über Autos":"Autos", 
                       "Berichte über Banken, Versicherungen und Geld":"Versicherungen/Geld", 
                       "Berichte über Bauprojekte in Ihrem Ort": "Bauprojekte vor Ort", 
                       "Berichte über Erziehungsfragen": "Erziehungsfragen", 
                       "Berichte über Hochschulen und Forschungsinstitute": "Hochschulen/Forschung", 
                       "Berichte über Kirchen": "Kirchen", 
                       "Berlichte über klassische Musik": "Klassische Musik",
                       "Berichte über Landespolitik": "Landespolitik",
                       "Berichte über lokale Schulen, Kindertageseinrichtungen":"Schulen/Kitas vor Ort", 
                       "Berichte über Messen": "Messen"}

umfrage_ressorts_dict ={"Kulturteil in der Zeitung":"Kulturteil", 
                        "Lokalteil in der Zeitung": "Lokalteil", 
                        "Politikteil in der Zeitung": "Politik", 
                        "Sportteil in der Zeitung": "Sportteil", 
                        "Wirtschaftsteil in der Zeitung": "Wirtschaft", 
                        "Wochenendbeilage": "Wochenendbeilage"}

umfrage_themen_dict = {"Berichte über Mode":"Mode", 
                       "Berichte über Partnerschaft und Beziehung": "Partnerschaft", 
                       "Berichte über Prominente und Klatsch": "Promis und Klatsch", 
                       "Berichte über Raumfahrt und Astronauten": "Raumfahrt", 
                       "Berichte über Umwelt und Klima": "Umwelt/Klima", 
                       "Berichte über Unglück und Verbrechen": "Unglück/Verbrechen", 
                       "Berichte über Wissenschaft und Forschung": "Wissenschaft", 
                       "Berichte über Renten": "Rente"}
                       
umfrage_themen_dict2 = {"Buchrezensionen":"Buchrezensionen", 
                        "Fernsehkritik und -tipps":"TV-Kritik und -tipps", 
                        "Humor":"Humor", 
                        "Informationen über Busse und Bahnen in Ihrem Ort": "ÖPNV vor Ort", 
                        "Informationen über einzelne Stadtteile": "einzelne Stadtteile", 
                        "Informationen über Essen und Trinken": "Essen und Trinken", 
                        "Informationen über Wellness und Kosmetik": "Wellness/Kosmetik"}

umfrage_themen_dict3 = {"Kino-Kritik": "Kino-Kritik", 
                        "Kreuzworträtsel/Sudoku":"Kreuzworträtsel/Sudoku", 
                        "Polizei- und Gerichtsberichte aus Ihrem Ort": "Kriminalität vor Ort", 
                        "Ratgeber und Verbraucherthemen": "Ratgeber", 
                        "Reportagen":"Reportagen",
                        "Theaterrezensionen":"Theaterrezensionen", 
                        "Veranstaltungshinweise":"Veranstaltungshinweise"}                     

#über diese Liste wird dann iteriert
umfrage_gesamt_liste = [umfrage_sport_dict,umfrage_gesell_dict, umfrage_ressorts_dict, 
                        umfrage_themen_dict,umfrage_themen_dict2, umfrage_themen_dict3]

# Zum Testen sollten bei einem neuen Kunden alle Dicts einmal mit der Update-Methode 
# zusammengezogen werden und mit den Inhalten der Kunden-Tabelle abgeglichen werden. 

# Die Fragen für die Pie-Charts
umfrage_piecharts = ["Seit wann haben Sie Ihre Zeitung abonniert?",
                     "Seit wann wohnen Sie am heutigen Wohnort?", 
                     "Zu welcher Tageszeit lesen Sie Ihre Zeitung?", 
                     "Wie groß ist Ihr Haushalt?", 
                     "Wie oft nehmen Sie durchschnittlich eine Ausgabe pro Tag in der Hand?", 
                     "Wie lange lesen Sie die Zeitung durchschnittlich pro Tag?", 
                     "Wie viele Seiten der Zeitung lesen Sie durchschnittlich?", 
                     "Wo lesen Sie in der Regel Ihre Zeitung?"]

umfrage_mini_bars_TV = ["Wie lange sehen Sie durchschnittlich an einem Tag fern?", 
                     "Wie oft sehen Sie Politik- und Wirtschaftsmagazine im TV?", 
                     "Wie oft sehen Sie Unterhaltungs-Shows im TV?", 
                     "Wie oft sehen Sie Spielfilme im TV?", 
                     "Welche Nachrichtensendungen sehen Sie sich hauptsächlich im TV an?", 
                     "Wie oft sehen Sie Nachrichten im TV?"]

umfrage_mini_bars_netz = ["Verfügen Sie über einen Internet-Zugang zu Hause?", 
                          "Verfügen Sie über ein internetfähiges Smartphone?", 
                          "Verfügen Sie über einen Internet-Zugang am Arbeitsplatz / Uni?", 
                          "Verfügen Sie über einen Tablet-PC?", 
                          "Wie lange surfen Sie durchschnittlich pro Tag im Internet?",
                          "Wie oft nutzen Sie das Internet?", 
                          ]

umfrage_demografie_mini_bars = ["Sie sind ...", "Wie alt sind Sie?", "Wie groß ist Ihr Haushalt?", 
                                "Wie hoch ist das monatliche Nettoeinkommen, das alle zusammen in Ihrem Haushalt haben - nach Abzug von Steuern und Sozialversicherung?"]

umfrage_demografie_bars = ["Welchen letzten Schulabschluss haben Sie?", "Sie sind zur Zeit ...", 
                           "Welcher Berufsgruppe ordnen Sie sich zu?"] 
                           


# ID-Nr für Downloads der Screenshots
id_nr = "1011" # Schwäbische Zeitung
      
#%% Function setlist() - Listen/Dicts verändern

# Die Funktion ersetzt/erweitert im Bedarfsfall bestehende Listen. 


def setlist(platzierung={}, darstellung={}, darstellungsf=[], 
            seitentitel={}, ressort=[], mantel = [], ressortdict={}, 
            sportarten=[], ausgaben=[], id_nummer="1009", seiten_lokal= [], 
            zeitung_attribute = []):
    
    # checken, ob die Parameter länger als 0 sind = neue Parameter
    # in dem Fall müssen die globalen Variablen verändert werden
    if len(platzierung)>0:
        print("Neue Platzierung_dict")
        global platzierung_dict
        platzierung_dict = platzierung
    
    if len(darstellung)>0:
        print("Neue Darstellungs_dict")
        global darstellung_dict
        darstellung_dict = platzierung_dict

    if len(darstellungsf)>0:
        print("Neue Liste Darstellungsform")
        global darstellungsform 
        darstellungsform = darstellungsf
        
    if len(seitentitel)>0:
        print("Neue Seitentitel gesetzt")
        global seitentitel_dict
        seitentitel_dict = seitentitel
    
    if len(ressort)>0:
        print("Neue Ressortliste gesetzt")
        global ressort_list
        ressort_list=ressort
    
    if len(mantel)>0: 
        print("Neue Mantelressorts gesetzt")
        global mantel_ressorts
        mantel_ressorts = mantel
    
    if len(ressortdict)>0:
        print("Neues Ressortdict gesetzt")
        global ressort_dict
        ressort_dict = ressortdict
    
    if len(sportarten)>0: 
        print("Neue Sportarten hinzugefügt")
        global liste_sportarten
        for elem in sportarten:
            liste_sportarten.append(elem)
    
    if len(ausgaben)>0: 
        print("Neue Ausgaben-Liste hinzugefügt")
        global ausgaben_liste
        for elem in ausgaben: 
            ausgaben_liste.append(elem)
        
    if id_nummer !="0": 
        print ("Neue ID für Downloads gesetzt")
        global id_nr
        id_nr = id_nummer
        
    if len(seiten_lokal) >0:
        print("Neue Liste für Lokalseiten gesetzt.")
        global seitentitel_lokal
        seitentitel_lokal = seiten_lokal
    
    if len(zeitung_attribute) >0: 
        print("Neue Liste für Zeitungsattribute in Umfrage-Sheets gesetzt.")
        global zeitung_attribute_dict
        zeitung_attribute_dict = zeitung_attribute
        
        
        
        
#%%
        
        
        
#%% ALLGEMEINES

#%% Entwicklung WErte im Messverlauf        
#Grafik Entwicklung - alle zwei Tage ein Balken + Gleitender Mittelwert
        # Unter target nimmt die Funktion die Argumente Lesewert, Blickwert, Durchlesewert entgegen. 
        # Unter title_text kann eine Überschrift für die Grafik übergeben werden
        
def grafik_entwicklung(prs, df, target="Lesewert", mean_line=0, legend="large", grid=True, title_text = False):
    if target == "Lesewert": 
        group_param = "Artikel-Lesewert (Erscheinung) in %"
    elif target == "Blickwert":  
        group_param = "Artikel-Blickwert (Erscheinung) in %"
    elif target == "Durchlesewert": 
        group_param == "Artikel-Durchlesewerte (Erscheinung) in %"
    df_ = df.groupby("Erscheinungsdatum", as_index=False).mean()
    
    # Grafik zeichnen
     #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
    
    # Werte für die Achsen werden festgelegt
    x = df_["Erscheinungsdatum"].apply(lambda x: x.strftime("%d.%m.%Y"))
    xn = range(len(x))
    labels = x
    y = df_[group_param]
    
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
    if grid==True:
        ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
     # Barcharts einzeichnen
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Lesewert")
            
    
                   
                   
        
    ax1.set_ylabel('Ø Lesewert in Prozent', color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=90, ha="center",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=20) # Labels werden ausgerichtet
    
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
    ax1.yaxis.label.set_size(22)
    
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    # Bei Bedarf Linie mit dem Durchschnitt einziehen
    if mean_line !=0:
        print(mean_line)
        labeltext = "Ø LW Seitentitel: {:1.1f}".format(float(mean_line)).replace(".", ",")
        linie = ax1.axhline(y=mean_line, xmin=0.01, xmax=0.99, color=set_font_color, label=labeltext)
    
    
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.44
    else:
        legend_height = 1.24
             
   
    
    if mean_line != 0:
        leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, linie], markerscale=140)
    else:
        leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars], markerscale=140)

    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    #plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    pos_left = 0.1 # 0.2
    pos_right=0.9 #0.8
    pos_top=0.90
    pos_bottom = 0.3
    
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_entwicklung_lesewert.png"
    plt.savefig(filename)
    
    # Prüfung ob Titel als Paramenter mitgegeben wurde, ansonsten automatisch eintragen
    if title_text == False: 
        title_text = "Entwicklung " + target 
    else: 
        title_text = title_text
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    picture_sheet(prs, filename, title_text=title_text)            
    
    return prs

      
        
        
#%% Initial-Funktion für Fragebogen 
        
        
        
def fragebogen(prs, df, liste = ausgaben_liste):
    print(df.head(2))
    
    # Identifizierung der Fragetexte ("Ihre Zeitung ist...")
    # geschieht über das dict zeitung_attribute_dict
    # die Anlage einer zusätzlichen Liste ist nicht notwendig, 
    # wir suchen über das dict
    
    # check ob attribute_df
    
    #Falls nötig: GESAMT
    analyse_mantel_abschluss(prs, df, liste_ressorts=mantel_ressorts)
    
    for elem in liste:
        
        elem_df = df[df["ZTG"]==elem]
        attribute_df = elem_df[elem_df["Fragetext"].isin(zeitung_attribute_dict)]
        #Berechnung der Werte, Erstellung des Diagramms, Schreiben auf PPTX
        #zeitungsattribute_berechnung(prs, attribute_df)
        #zeitung_themen(prs, elem_df)
        #deckblatt(prs, "Umfrage" + elem)
        #umfrage_pie(prs, elem_df)
        
        #mini_bars(prs, elem_df, title_text = "Zeitungsnutzung", liste=umfrage_mini_bars_TV)
        #mini_bars(prs, elem_df, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz)
        
        
        
    return prs
        



#%% Berechnung Attribute_Werte der Zeitung  
# Die Funktion erstellt mehrfarbige, gestapelte  Bar-Charts für die Darstellung
# der Attribut auf den Satz "Ihre Zeitung ist..." 
# Als Parameter benötigt die Funktion das pptx-Objekt und die
    # Umfragedaten

def zeitungsattribute_berechnung(prs, df):
    
    print("Funktion Zeitungsattribute_berechnung gestartet")
    df = df[["Welle", "Fragetext", "Antworttext"]]
    df_group = df.groupby(["Fragetext", "Antworttext"], as_index=False).count() 
    # DAten von long auf wide umbauen, um Grafik leichter zeichnen zu können
    df_group = df_group.pivot(index= "Fragetext", columns="Antworttext", 
                              values= "Welle")
    
    df_group.reset_index(inplace=True)
    
    df_group.fillna(0, inplace=True)
    
    df_pivot = df_group.reindex(columns=["Fragetext", "trifft nicht zu", "trifft eher nicht zu", "trifft eher zu", "trifft voll zu"])
    
    
    
    # START PLOTTING
    # Variablen festlegen
    x = df_pivot["Fragetext"]
    df_pivot["shortnames"] = df_pivot["Fragetext"].apply(lambda x: zeitung_attribute_dict[x])
    labels = df_pivot["shortnames"]
    xn = range(len(x))
    width = 0.3
    #Schriftfarbe und Farbe der Ticks festlegen, fig/ax aufrufen
    set_font_color ="#8c8f91" 
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(18,8))
    #setzt Schriftart und Schriftgrößé für linke Y-Achse
    plt.yticks(fontproperties=campton_light)
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
   
    p1 = ax1.bar(xn, df_pivot["trifft nicht zu"], width, 
                 label=("trifft nicht zu"), color="#f9a424")
    p2 = ax1.bar(xn, df_pivot["trifft eher nicht zu"], width, \
             bottom=df_pivot["trifft nicht zu"], label=("trifft eher nicht zu"), color="#fce9cc")
    p3 = ax1.bar(xn, df_pivot["trifft eher zu"],width,\
             bottom=df_pivot["trifft nicht zu"]+
             df_pivot["trifft eher nicht zu"], label=("trifft eher zu"), 
             color="#f9d69f")
    p4 = ax1.bar(xn, df_pivot["trifft voll zu"], width, \
             bottom=df_pivot["trifft nicht zu"]+df_pivot["trifft eher nicht zu"]
             +df_pivot["trifft eher zu"],color="#fcc267", 
             label=("trifft voll zu"))
    
    #Labels und Ticks festlegen
    ax1.set_ylabel('', color= set_font_color, \
                   fontproperties=campton_light, fontsize=30)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xlabel("")
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    # Legende
    
    handles,labels = ax1.get_legend_handles_labels()
    handles = [handles[3], handles[2], handles[1], handles[0]]
    labels = [labels[3], labels[2], labels[1], labels[0]]
    leg = ax1.legend(handles, labels, bbox_to_anchor=(1.05, 1), loc=2, borderaxespad=0., markerscale=140)
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=30)
        
    #Lage der Bars / Sichtbarkeit der Spines  
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
    #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    #Grid-Linien
    ax1.grid(color= "#e0e0e0", linestyle='-', linewidth=1, axis="y")
          
    # Titel hinzufügen
    figure_title= "Ihre Zeitung ist..."
    ax1.text(0.5, 1, figure_title, horizontalalignment='center',
             verticalalignment='baseline', 
             transform = ax1.transAxes,fontproperties=campton_light,
             color = set_font_color, fontsize=50, 
             bbox=dict(facecolor='none', edgecolor='none', pad=30) )  
            #[ 'center' | 'top' | 'bottom' | 'baseline' ]
#    
#if label_position == "large": 
#       p_left = 0.12
#       p_right = 0.92
#       p_top = 0.85
#       p_bottom = 0.4
#   
#   if label_position == "xlarge": 
#       p_left = 0.14
#       p_right = 0.92
#       p_top = 0.85
#       p_bottom = 0.45

    plt.subplots_adjust(left=0.12, right=0.92, top=0.95, 
                        bottom=0.4)   
    plt.margins(0.05,0.19)
    plt.tight_layout()
    #plt.tight_layout(pad=4, w_pad=0, h_pad= 1.0)
    # Canvaseinstellung / Position des Plots
    #
   
    
    
    filename = "grafik_ztg_attribute.png"
    plt.savefig(filename, bbox_inches="tight")
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    # Vergabe des Titeltextes
    title_text = "Zeitungsnutzung"
    picture_sheet(prs, filename, title_text=title_text)         
    #savefig("filename.pdf", bbox_inches = 'tight',
    #pad_inches = 0)
    
    #gca().xaxis.set_major_locator(NullLocator())
    #gca().yaxis.set_major_locator(NullLocator())
    return prs     
        
        


#%% Welche Themen - Balkengrafiken quer für Fragebogen
    # Die Funktion benötigt das PPTX-Objekt und das Dataframe der Umfrage. 
    # Erstellt werden Querbalkengrafiken für die Frage: "Welche der folgenden Themen 
    # interessieren Sie besonders?"
    # Die Antworten,  die jeweils auf eine Folie sollen, sind unter Listen und Objekten
    # jeweils in Listen angegeben. Die müssen bei Bedarf geändert werden. 
    
def zeitung_themen(prs, df):
    fragetext = "Welche der folgenden Themen interessieren Sie besonders in der Zeitung?"
    df_ = df[df["Fragetext"]==fragetext]
    
    for elem in umfrage_gesamt_liste: 
        
        df_loop = df_[df_["Antworttext"].isin(elem)]
        
        df_group = df_loop.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle")
        
        
        # Grafik zeichnen
     
        df_group["Antwort_short"] = df_group["Antworttext"].apply(lambda x: elem[x])
        answer = df_group["Antwort_short"]
        y_pos = range(len(answer))
        nennung = df_group["Welle"]
        bar_color = "#f77801"
       
        width = 0.3
        set_font_color = "#8c8f91"
    
    
        #Schriftfarbe und Farbe der Ticks festlegen
        
        sns.set_style("white")
        fig, ax1 = plt.subplots(figsize=(19,7))
        #setzt Schriftart und Schriftgrößé für linke Y-Achse
        plt.yticks(fontproperties=campton_light)
        plt.xticks(fontproperties=campton_light)
        ax1.tick_params(axis='x', labelsize=25, colors= set_font_color)
        ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
        ax1.set_yticks(y_pos)
        ax1.set_yticklabels(answer)
    
        ax1.barh(y_pos, nennung, color=bar_color, align="center")
        
        plt.margins(x=0.05) # ziehen Bars näher an die Achse
             #obere Linie ausblenden
        ax1.spines["top"].set_visible(False)
            #ax1.spines["left"].set_color("gray")
        ax1.spines["top"].set_visible(False)
        ax1.spines["bottom"].set_visible(False)
        ax1.spines["left"].set_visible(False)
        ax1.spines["right"].set_visible(False)
    
    
    
        #Grid-Linien
        ax1.grid(color= "#e0e0e0", linestyle='-', linewidth=1, axis="x")
    
        # Titel hinzufügen
        figure_title= fragetext
        plt.text(0, 1.0, figure_title,
                 #horizontalalignment='center',
    
                transform = ax1.transAxes, 
                fontproperties=campton_light, 
                color = set_font_color, 
                fontsize=30)
        
        
        #plt.subplots_adjust(left=0.12, right=0.92, top=0.85, 
                            #bottom=0.4)   
       # plt.margins(0.05,0.19)
        plt.tight_layout()
        
        
        filename = "grafik_ztg_themen .png"
        plt.savefig(filename, bbox_inches="tight")
        plt.close()
        # Plot wird auf PPTX-Sheet gezogen
        # Vergabe des Titeltextes
        title_text = "Zeitungsnutzung"
        picture_sheet(prs, filename, title_text=title_text)  
    return prs


#%% Pie-Charts (klein) Zeitungsnutzung
    # Diese Funktion schreibt jeweils zwei Pie-Charts auf ein PPTX-Dokument. 
    # Die Reihenfolge der Pie-Charts ist unter Listen und Objekte in der Liste 
    # umfrage_piecharts festgelegt. 
def umfrage_pie(prs, df):
    print("Funktion umfrage_pie gestartet")
    counter=0
    for elem in umfrage_piecharts:
        counter += 1
        print(elem)
        df_ = df[df["Fragetext"]==elem]
        
        df_group = df_.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle", ascending= False)
        
        
        
        sizes = df_group["Welle"]
        set_font_color ="#8c8f91"
        print("size: ")
        print(sizes)
        # legt Position der Zahlen innerhalb der Pie-Chart fest
        def autodistance(autopct):
            print(autopct)
            if autopct >=5:
            
                return 1.15 # 1.2 bei kleinen Zahlen
            else:
                return 0.8
        
        # Erstellt die tatsächlichen Nutzerzahlen anhand der automatisch generierten Prozentzahlen
        def make_label(sizes):
            total = sizes.sum()
            def my_label(pct):
                return round(total*(pct/100)).astype(int)
            mylabel = my_label
            
            #if mylabel >=6:
               # pct_distance = 0.8
            #if mylabel <6:
               # pct_distance = 1.2
            return mylabel
        
        # Falls nötig: Explode wird für jedes Element errechnet
        explode = ()
        for i in range(len(sizes)):
            explode = explode +(0.03,)
        print (explode)
        colors =["#f9a424", "#fce9cc", "#f9d69f", "#fcc267"]
        fig, ax1 = plt.subplots(figsize=(8,8))
        test = make_label(sizes)
        print(test)
        print("--------------------------")
        ax1.pie(sizes, radius=1, frame=False, shadow=False, autopct=make_label(sizes), startangle=90, colors=colors,\
                textprops={'fontsize': 20, "color": set_font_color},pctdistance =0.8, 
                wedgeprops={"linewidth":2, "edgecolor":"white"})
        ax1.axis("equal")
        
            
        
        handles, labels = ax1.get_legend_handles_labels()
        labels = df_group["Antworttext"].values
       
        
        leg = ax1.legend(handles, labels=labels, #bbox_to_anchor=(1, 0.2), 
                         markerscale=14, mode="expand", 
                         borderaxespad=0., loc=8) #1.Wert x, 2. Wert y
        for text in leg.get_texts(): 
            plt.setp(text, color= set_font_color, size=20)
        
        # Überschrift
        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
        figure_title= elem
        posx=0.5
        posy=0.82
        if len(figure_title) >40:
            figure_title = "\n".join(textwrap.wrap(elem, width=35))
            posy=0.8
        print(figure_title)    
        
        plt.text(posx, posy, figure_title,
             horizontalalignment='center',
             
             transform = ax1.transAxes, 
            fontproperties=campton_light, 
            color = set_font_color, 
            fontsize=28)
        
        
        filename = "grafik_pie.png"
        filename2 = "grafik_pie2.png"
        
        #plt.tight_layout()
        #plt.savefig(filename)
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    
        
        pos_left = 0.2 # 0.2
        pos_right=0.7 #0.8
        pos_top=0.98
        pos_bottom = 0.1
        plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                       bottom=pos_bottom)
        
        
        # plt.savefig(filename, bbox_inches="tight")
        title_text = "Zeitungsnutzung"
        if counter%2!=0:
            plt.savefig(filename)
        elif counter%2==0:
            plt.savefig(filename2)
            double_picture_sheet(prs, filename, filename2, title_text=title_text)
        plt.close()
        
        #plt.savefig(filename, bbox_inches="tight")
        
        #picture_sheet(prs, filename, title_text=title_text) 
        #double_picture_sheet(prs, filename, title_text=title_text)
        
    return prs





#%% Fragebogen Mini-Barchart
    # Die Fuktion legt die Fragen fest, die per Mini-Barchart auf PPTX-Docs gezogen
    # werden (je zwei auf eine Seite). 
    # Die Funktion benötigt eine Liste,  in der die Fragen angegeben sind, die 
    # genutzt werden sollen. Voreingestellt ist umfrage_mini_bars_TV, es gibt auch die 
    # liste umfrage_mini_bars_netz. Die Listen finden sich unter Listen und Objekte. 
    
    # Die Überschrift kann mit title_text individuell eingestellt werden. 
   
        
def mini_bars(prs, df, title_text = "Zeitungsnutzung", liste=umfrage_mini_bars_TV):
    counter = 0
    print("Funktion mini_bars gestartet")
    if len(liste)==0:
        print("Achtung, Funktion mini_bars ohne Liste mit Umfrage-Themen")
    for elem in liste:
        counter +=1
        df_ = df_ = df[df["Fragetext"]==elem]
        df_group = df_.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle", ascending= False)
        
        set_font_color ="#8c8f91"
        
        x = df_group["Antworttext"]
        
            
        # Größe der X-Ticks, wird bei längeren Texten Internetnutzung geändert    
        ticksize=30  
        
        def check_labels(elem):
            if len(elem)>19:
                return elem[:16] + "..."
            else:
                return elem
            
        def check_labels_internet(elem):
            if len(elem)> 8:
                return_label = "\n".join(textwrap.wrap(elem, width=8))
                ticksize = 20
                return return_label
            else:
                return elem
          
        ticksize = 30
        if title_text == "Zeitungsnutzung":  
            labels = x.apply(check_labels)
            rotation=45
            ha="right"
            # Canvaseinstellung / Position des Plots
            # Function nutzt Voreinstellung aus Parametern
            p_left = 0.3
            p_right = 0.8
            p_top = 0.7
            p_bottom = 0.45
            figsize=(10,8)
            
            
        elif title_text == "Internetnutzung": 
            labels = x.apply(check_labels_internet)
            
            print(labels)
            rotation= 0
            ha="center"
            # Canvaseinstellung / Position des Plots
            # Function nutzt Voreinstellung aus Parametern
            p_left = 0.2
            p_right = 0.8
            p_top = 0.6
            p_bottom = 0.2
            figsize = (10,8)
            if len(labels)>9:
                ticksize=20
            elif len(labels)<=9:
                ticksize=30
        #TODO ticksize muss sich verkleinern, wenn der Text der X-Ticks länger wird    
        xn = range(len(x))
        y = df_group["Welle"]
        
         # Seaborn-Style und Größe des Plots festlegen
        sns.set_style("white")
        fig, ax1 = plt.subplots(figsize=figsize)
    
        #setzt die linke Y-Achse in Campton light
        # rechte Y-Achse können wir erst zum Code-Ende ansteuern
        #plt.yticks(fontproperties=campton_light)
        
        #Grid einbauen
        #ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
        
        bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Anzahl Antworten")
        
        
        ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
        
        
        ax1.set_xticklabels(labels = labels, rotation=rotation, ha=ha,  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=ticksize) # Labels werden ausgerichtet
        
        for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
            height = p.get_height()
        
            txt = '{:1.0f}'.format(height).replace(".", ",")
            ax1.text(p.get_x()+p.get_width()/2., height + 5, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 30, weight = 1000)
        
        #ax1.yaxis.label.set_size(22)
        plt.margins(x=0.03) # ziehen Bars näher an die Achse
        
        ax1.spines["top"].set_visible(False)
        #ax1.spines["left"].set_color("gray")
        ax1.spines["top"].set_visible(False)
        ax1.spines["bottom"].set_visible(False)
        ax1.spines["left"].set_visible(False)
        ax1.spines["right"].set_visible(False)
        
        
         # Überschrift
        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
        figure_title= elem
        posx=0.5
        posy=1.5
        if len(figure_title) >40:
            figure_title = "\n".join(textwrap.wrap(elem, width=35))
            
        print(figure_title)    
        
        plt.text(posx, posy, figure_title,
             horizontalalignment='center',
             
             transform = ax1.transAxes, 
            fontproperties=campton_light, 
            color = set_font_color, 
            fontsize=28)
        
        ax1.axes.get_yaxis().set_visible(False)
        
        #plt.tight_layout()
    
       
        
        plt.subplots_adjust(left=p_left, right=p_right, top=p_top, 
                        bottom=p_bottom)
        
        filename1 = "grafik_lesewert_mini_bars.png"
        filename2 = "grafik_lesewert_mini_bars2.png"
        title_text = title_text
        
        
        if counter%2 !=0:
            
            plt.savefig(filename1)
        elif counter%2 ==0:
            
            plt.savefig(filename2)
            double_picture_sheet(prs, filename1, filename2, title_text=title_text)
        plt.close()
        
    return prs
        
#%%
    

#%% ALLGEMEINE ERKENNTNISSE
    

#%% Initialfunktion Erkenntnisse

#%% Lesetage 
    # Funktion erstellt Barchart über die Lesetage und die durchschnittliche Zahl von Lesern 
    # an den jeweiligen Tagen





#%% Lesezeiten    
''' Diese Helper-Function erstellt Analysen über den Lesezeitraum (Lesezeitpuinkt am Tag. Sie benötigt 
das Dataframe ScanAuswertungen aus der Lesewert-Datenbank. Ansonsten kümmert sie 
sich um sich selbst. 

'''

def lesezeit(prs, df_scans):
    # TODO: Mal mit Martin klären, ob die nicht Treffer "Falsch" auch gewertet werden 
    df_time = df_scans[df_scans["Treffer"]=="WAHR"]
    print ("df_time---------------")
    print(df_time.count())
    # Zusammenführen der Datumsangaben
    
    df_time["time"] = df_time["Erfassungsdatum"] + "-" + df_time["Erfassungsuhrzeit"]
    
    #Erstelle Datetime-Format
    df_time["time"] = pd.to_datetime(df_time["time"], format="%d.%m.%Y-%H:%M:%S")
    
    # Filtere Jahre, nur 2018 erlaubt
    mask = df_time["time"].dt.year == 2018
    df_time = df_time[mask]

    # Erstellen einer Spalte mit den Lesestunden
    df_time["hour"] = df_time["time"].dt.hour
    
    
    # neues Dataframe mit 24 Columns, um für jeden einzelnen Erscheinungstag 
    # zu zählen
    df_new = pd.DataFrame(columns = np.arange(24))
    list_days = df_time["Erfassungsdatum"].unique()
    
    for elem in list_days:
        df_ = df_time[df_time["Erfassungsdatum"]==elem]
        daily_list = []
        for i in range (24):
            reader = df_[df_["hour"]==i].WellenteilnahmenId.nunique()
            daily_list.append(reader)
        df_new = df_new.append(pd.Series(daily_list, index=df_new.columns), 
                               ignore_index= True)
    
    # neues Dataframe um die Reader der einzelnen Tage pro Stunde zusammenzurechnen
    col_names = ["hour", "reader"] 
    df_sum = pd.DataFrame(columns=col_names)
    
    for i in range(24):
        hour = i
        reader = df_new[i].sum()
        #df_sum = df_sum.append(pd.Series([hour,reader], index=df_sum.columns), ignore_index= True)
        
        df_sum = df_sum.append(pd.Series([hour,reader], index=df_sum.columns), ignore_index= True)
    
    # converting df_sum-columns into int
    df_sum["reader"] = pd.to_numeric(df_sum["reader"]).astype(float)
    
    
    
   # Anzahl Erscheinungstage: 
    ET_tage = df_time["Erfassungsdatum"].nunique()
    
    
    x= df_sum["hour"]
    x_ = np.array(x)
    labels = x
    xn = range(len(x))
    y = df_sum["reader"] /ET_tage
    y_ = np.array(y)
    yn = df_sum["reader"] /ET_tage
    
    x_new = np.linspace(x_.min(), x_.max(),500)
    f = interp1d(x_, y_)
    y_smooth=f(x_new)
    
    
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)
    
    #ax1.plot(x_new,y_smooth, color="#f77801", label="Lesewert")
    ax1.set_ylabel('Ø Leserzahl', color="#aeb0b2", fontproperties=campton_light, fontsize=50)
    ax1.set_xlabel('Uhrzeit', color="#aeb0b2", fontproperties=campton_light, fontsize=20)
    ax1.xaxis.set(ticks=range(0, len(xn))) # wir müssen die Anzahl der Ticks händisch festlegen
    ax1.set_xticklabels(labels = x, rotation=45, ha="right",  weight=800, color="#aeb0b2", \
                        fontproperties=campton_light, fontsize=30) # Labels werden ausgerichtet
                        #fontname=campton_light) # Labels werden ausgerichtet  
    fill = ax1.fill_between(xn, yn, alpha=1, color="#f77801",label="Ø Artikel/Tag") # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
    
    ax1.yaxis.label.set_size(22)
    
    ax1.spines["top"].set_visible(False)
    ax1.spines["left"].set_color("gray")
    ax1.spines["bottom"].set_color("gray")
    ax1.spines["right"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors="#aeb0b2", )
    
    
    plt.tight_layout()
    plt.subplots_adjust(left=0.08, right=0.92, top=0.85, bottom=0.3)
    filename = "grafik_zeiten3.jpg"
    plt.savefig(filename)
    
    picture_sheet(prs, filename, title_text= "Wann wird gelesen?")
    
    return prs




#%% LW-Marken - Textanzeige der Daten
    
''' Diese Funktion wirft lediglich die Werte als Text heraus. Funktioniert
nur für die Schwäbische, muss für jede Zeitung neu gemacht werden. 
'''

def marken_analyse(df, df_doublesplitid=False, df_nichtkum=False):
    # Werte Gesamt
    print("WERTE GESAMT")
    gesamt_lw = df['Artikel-Lesewert (Erscheinung) in %'].mean()
    gesamt_bw = df['Artikel-Blickwert (Erscheinung) in %'].mean()
    gesamt_dw = df['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    print ("LW gesamt: {:1.1f}%".format(gesamt_lw).replace(".", ","))
    print ("BW gesamt: {:1.1f}%".format(gesamt_bw).replace(".", ","))
    print ("DW gesamt: {:1.1f}%".format(gesamt_dw).replace(".", ","))
    
    # Werte Lokales
    #df_l = df[(df["Ressortbeschreibung"]=="Lokales") | (df["Ressortbeschreibung"]=="Lokalsport")]
    print("")
    print("")
    print("WERTE LOKALTEILE")
    # check ob df_doublesplitid vorhanden ist
    if isinstance(df_doublesplitid, pd.DataFrame):
        print("Datensatz mit doppelten SplitId gefunden...")
        df_lokal = df_doublesplitid.copy()
        
    else: 
        print("Kein Datensatz mit doppelten SplitId gefunden")
        print("Analyse mit kumulierten Daten.")
        df_lokal = df.copy()
    
    df_l = df_lokal[df_lokal["Ressortbeschreibung"]== "Lokales"]
    
    ausgaben = df_l["ZTG"].unique()
    print (ausgaben)
    for ausg in ausgaben: 
        df_ = df_l[df_l["ZTG"]==ausg]
        lw = df_['Artikel-Lesewert (Erscheinung) in %'].mean()
        bw = df_['Artikel-Blickwert (Erscheinung) in %'].mean()
        dw = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
        print ("LW Lokalteil " + ausg +": {:1.1f}%".format(lw).replace(".", ","))
        print ("BW Lokalteil " + ausg +": {:1.1f}%".format(bw).replace(".", ","))
        print ("DW Lokalteil " + ausg +": {:1.1f}%".format(dw).replace(".", ","))
    print("-------------------------end-------------------")
   
    
    
    # Mantel
     
 
    print("WERTE MANTEL")
    mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
    
    mantel_lw = mantel['Artikel-Lesewert (Erscheinung) in %'].mean()
    mantel_bw = mantel['Artikel-Blickwert (Erscheinung) in %'].mean()
    mantel_dw = mantel['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    print ("LW Mantel: {:1.1f}%".format(mantel_lw).replace(".", ","))
    print ("BW Mantel: {:1.1f}%".format(mantel_bw).replace(".", ","))
    print ("DW Mantel: {:1.1f}%".format(mantel_dw).replace(".", ","))
    print("-------------------------end-------------------")
    
    
    
    print("WERTE EINZELAUSGABEN (NICHTKUM)")
    
    # Check ob nichtkumulierte Werte übergeben worden sind. 
    # Ansonsten nehmen wir die kumulierten
    # Check über Type und isinstance()
    print(df_nichtkum.head(2))
    if isinstance(df_nichtkum, pd.DataFrame):
        print("Eigener Datensatz für nichtkumulierte Werte gefunden...")
        df_total = df_nichtkum.copy()
        
    else: 
        print("Keine nichtkumulierten Daten vorhanden.")
        print("Analyse mit kumulierten Daten.")
        df_total = df.copy()
        
        
        
        
    for ausg in ausgaben: 
        df_ = df_total[df_total["ZTG"]==ausg]
        lw = df_['Artikel-Lesewert (Erscheinung) in %'].mean()
        bw = df_['Artikel-Blickwert (Erscheinung) in %'].mean()
        dw = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
        
        print("LW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(lw).replace(".", ","))
        print("BW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(bw).replace(".", ","))
        print("DW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(dw).replace(".", ","))
     




#%% Analyse Ressrots/Seitentitel nach Lesewert
    
'''
Diese Helperfunction wirft eine Grafik zum Lesewert und der Erscheinungshäufigkeit
einzelner Ressorts oder Seitentitel heraus. 


Target="ressort" ist die Voreinstellung. Sollen Seitentitel analysiert werden, 
target="seitentitel"

Die Funktion benötigt für Target="ressort" nur den normalen, bereinigte Datensatz. 
Für Target = "seitentitel" benötigen wir den Datensatz auf Ressortebene. 

special: Hier fasse ich verschiedene Anwendungsmöglichkeiten zusammen. 
1 = Grafik mit einer Y-Achse links, LW für einzelne Ressorts, Linie mit allg.DS
2 = Grafik mit zwei Y-Achsen, Artikelhäufigkeit pro Ressort (absolute Zahl)



minimize, Standard 5, zeigt an ab welcher Zahl von Artikeln ein Ressort/Seitentitel
in die Auswertung aufgenommen wird.

label_position gibt an, wie viel Platz die Labels unten und links bekommen. 
Voreingestellt auf normal. Auch möglich:  large und xlarge für besonders lange 
Labelnamen. strange sind spezialformate

Über die ressort_liste kann der Funktion noch explizit mitgeteilt werden, 
welche Ressorts/Seitentitel in die Analyse einfließen sollen

Unter Special können Spezialfälle für einzelne Zeitungstitel angelegt werden. 

order = "ok" meint, dass wir an der Sortierung der Ressorts nach Seitenzahl /Median
nicht mehr händisch eingreifen müssen. order = "new" heißt, dass wir manuel eingreifen
'''



def grafik_lesewert(prs, df, target="ressort", minimize=5, label_position="normal",
                    ressort_liste=[], special=False, title_text="", order="ok", legend="normal", sort="Seitennummer", 
                    article="total"): 
    
   
  
   
   group_param=""
   
   #Festlegung welche Werte hier gesucht werden
   
   if target=="ressort":
       group_param = "Ressortbeschreibung"
       
   elif target == "Darstellungsform": 
       group_param = "Darstellungsform"
       
       
   elif target == "seitentitel":
       group_param = "Seitentitel"
       
     
   # Abfrage, ob eine Auswahl-Liste mitgegeben wurde. Falls ja: DF bereinigen.      
   if len(ressort_liste)>0:
       
       df = df[df[group_param].isin(ressort_liste)]    
       
       
   
   df_group = df.groupby(group_param, as_index=False)\
       .agg({"SplitId":"size", "Seitennummer":"median", "Artikel-Lesewert (Erscheinung) in %":"mean"})
   sort_value = sort
   df_group = df_group[df_group["SplitId"]>=minimize]
   if sort == "Seitennummer": 
       #niedrigste Seitennummer nach vorne
       df_group = df_group.sort_values(by=sort)
   else:
       # hnöchster Lesewert nach vorne
       df_group = df_group.sort_values(by=sort, ascending=False)
   
   
   # Lokalteil identifizieren
   
   
       
   # Reihenfolge ändern, damit Lok 1 vorne steht:
   if order == "new": 
       list_lok = df.Seitentitel.unique()
       if "Biberach" in list_lok:
           
           df_group = df_group.reindex([4,2,1,5,8,7,9, 6, 3])
           
       elif "Friedrichshafen" in list_lok:
           print("------------- TEST FHA -------------------------")
           print (df_group)
           
           
           df_group = df_group.reindex([4,1,3,6,5,9,8,2])
           print (df_group)
   
  
   df_group.dropna(inplace=True)
   df_group.reset_index(drop=True, inplace=True)
  
    # STEP 1 - Berechnung des durchschnittlichen Lesewerts aus den Einzelwerten
    # der Ressorts(Seitentitel)
   
   lw_durchschnitt = df_group["Artikel-Lesewert (Erscheinung) in %"].mean()
   
  
   # SPECIAL für Schwäbisch
   # Werte für Titelseite und Lokales werden gedrittelt
   # Grund: Da bei Titel und Lokalem jeweils alle drei 
   
   
   # TODO Viel zu gefärhlich... man muss die Index-Nummer der jeweiligen 
   # Werte angeben... totale Fehlerquelle
   
   if special==True: 
       #df_group[df_group["Ressortbeschreibung"]=="Lokales"]["]
       mask = df_group["Ressortbeschreibung"]=="Lokales"
       val1 = df_group[mask].SplitId
       val_lok = val1 / 3
       
       df_group.set_value(10, "SplitId", val_lok)
       
       mask2 = df_group["Ressortbeschreibung"]=="Titel"
       val2 = df_group[mask2].SplitId
       val_titel = val2 / 3
       df_group.set_value(0, "SplitId", val_titel)
       
       mask3 = df_group["Ressortbeschreibung"]=="Lokalsport"
       val3 = df_group[mask3].SplitId
       val_loksport = val3/3
       df_group.set_value(11, "SplitId", val_loksport)
       
       
      
     
    
    
    
    
    
    # Werte für Achsen festlegen
   #Achsen-Werte festlegen
   x = df_group[group_param]
   
   if target == "ressort":
       
         # schiebt die beiden ersten Lokalseiten im Mantel ans Ende der lokalen 
       # SEiten, vor allem hinter die Lok 1
   
    
        df_group["shortnames"] = df_group[group_param].apply(lambda x: ressort_dict[x])
       
        labels = df_group["shortnames"]
        
   
    
    
   
    
   # Seitentitel sind oft zu lang, daher ersetzen wir die Labels mit Kurzversionen    
   if target == "seitentitel":
       print (df_group[group_param])
       
       df_group["shortnames"] = df_group[group_param].apply(lambda x: seitentitel_dict[x])
       labels = df_group["shortnames"]
   
   
    
   xn = range(len(x))
   y = df_group["Artikel-Lesewert (Erscheinung) in %"]
   et = df.Erscheinungsdatum.nunique()
   yn = df_group["SplitId"] / et
   
   
   # Titel der Grafik einstellen
   
   
   if target=="ressort": 
       title_text = "Lesewert nach Ressorts"
       

       
   elif target == "seitentitel":
       title_text = "Lesewert nach Seitentiteln"
       
   # Postion und Lage der Labels auf der X-Achse bestimmen    
   
       
   p_left = 0.08
   p_right = 0.92
   p_top = 0.85
   p_bottom = 0.33
   
   if label_position == "large": 
       p_left = 0.12
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.4
   
   if label_position == "xlarge": 
       p_left = 0.14
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.48
       
   if label_position == "strange":
       p_left = 0.14
       p_right = 0.80
       p_top = 0.75
       p_bottom = 0.45
    
   
   # Grafik erstellen
   plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
                  pos_left = p_left, pos_right=p_right, pos_top=p_top, \
                  pos_bottom = p_bottom, article = article, grid=False,\
                  title_text=title_text, axis=1, mean_line=0, legend=legend)
    
   
   # STEP 2 - Berechnung durchschnittlicher Lesewert plus Häufigkeit je Ressort
   # Grafik erstellen
   plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn=yn, \
             pos_left = p_left, pos_right=p_right, pos_top=p_top, \
             pos_bottom = p_bottom, article = article, grid=False,\
             title_text=title_text, axis=2, mean_line = 0, legend=legend)
   
   
   
   
   return prs 
 
    
#%% Ressorts nach Geschlecht
   # TODO: Daten fehlen noch

#%% Wochentage nach Lesewert und durchschnittlicher Artikelanzahl pro Tag
#Func Plot von Lesewert nach Erscheiungstag
'''
Lesewert nach Erscheinungstag

Diese Funktion  übernimmt einen vorgefilterten Datensatz (zum Beispiel alle 
Artikel eines Ressorts etc.) und zeigt dann alle Erscheinungstage mit 
durchschnittlichem Lesewert und durchschnittliche Anzahl der erschienen Artikel 
an diesem Tag. 

Optional kann eine neue Überschrift vergeben werden. 



'''

def lesewert_erscheinung(df, prs, title_text="Lesewert nach Wochentagen"):
    
    df_ = df.copy()
    erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    df_["Tag"] = df_["Erscheinungsdatum"].apply(lambda x: x.strftime('%A'))
    
    df_ = df_[["SplitId", "Ressortbeschreibung", "Erscheinungsdatum", "weekday", "Tag", \
               "Artikel-Lesewert (Erscheinung) in %", "Seitennummer", "Darstellungsformen"]]
    
    
    df_number = df_.groupby(["Tag", "Erscheinungsdatum"], as_index=False).count()
    df_2 = df_number.groupby("Tag", as_index = False).count()
    df_2 = df_2.rename(columns={"Erscheinungsdatum": "Tageszahl"})
    df_2 = df_2[["Tag", "Tageszahl"]]
    # Groupby um Anzahl der Artikel und durchschnittliche Lesewerte zu erhalten
    df_group = df_.groupby("Tag").agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean", \
                                           "Erscheinungsdatum":"count"})
        
    df_group = df_group.reindex(["Montag", "Dienstag", "Mittwoch", "Donnerstag", "Freitag", "Samstag"])
    
    # Columns ordnen
    df_group = df_group.rename(columns={"SplitId":"Artikelanzahl", "Artikel-Lesewert (Erscheinung) in %":"LW"})
    df_group['Tag'] = df_group.index
    
    #Verschmelzung der beiden Tabellen
    df_group = pd.merge(df_group, df_2, on="Tag", how="left")
    df_group.dropna(inplace=True)
       
    #Achsen-Werte festlegen
    x = df_group["Tag"]
    labels = x
    xn = range(len(x))
    y = df_group["LW"]
    yn = df_group["Artikelanzahl"]  / df_group["Tageszahl"]
    
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.33, article = "mean", grid=False,\
                  title_text = title_text, axis=2, legend="strange")


    return prs






   
#%% Lesewert/Artikelzahl nach Darstellungsformen

''' 
Die Funktion 
darstellungsformen()
zeichnet eine Grafik mit zwei y-Achsen zu den Darstellungsformen und der Anzahl 
der Artikel in diesen Formen. 

Der Parameter minimum sagt aus wie häufig eine Darstellungsform vorkommen muss, um 
gemessen zu werden. 
By default steht der Wert auf 5. 

Ansonsten muss der User nur das entsprechende Dataframe eingeben, das 
ausgewertet werden soll. 

'''

def darstellungsformen(prs, df, minimum = 5, ):
    # liste heißt Darstellu´ngsform
    mask = df["Darstellungsformen"].notnull()
    df = df[mask]
    
    #Diese Funktion spaltet die Darstellungsformen-Zellen auf
    # außerdem werden hier die beiden Darstellungsformen  BB und BN zu BI zusammengefasst
    def splititup(elem):
        splitter = [x.strip() for x in elem.split(',')]
        
        for x in splitter: 
            if x in darstellungsform: 
                if x == "BB" or x == "BN":
                    return "BI"
                else: 
                    return x
            else:
                pass

    #neue Spalte DF wird angelegt, hier werden erwünschte Darstellungsformen verzeichnet
    df["DF"] = df["Darstellungsformen"].apply(splititup) 

    # Tabelle wird verkleinert, Zeilen ohne Eintrag in DF werden gelöscht
    df_ = df[["SplitId", "Erscheinungsdatum", "Darstellungsformen", "DF", "Artikel-Lesewert (Erscheinung) in %"]]
    mask_df = df_["DF"].notnull()
    df_ = df_[mask_df]

    #Anzahl Erscheinungstage wird ermittelt
    erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    # Tabelle mit Darstellungsform, Anzahl Artikel und Durchschnittslesewert wird erstellt
    df_DF = df_.groupby(["DF"], as_index=False).agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})\
    .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
    
    
    df_DF = df_DF[df_DF["SplitId"]>=minimum]
    # GRAFIK ANFERTIGEN
    
    # Achsen festlegen
    x= df_DF["DF"]
    labels = df_DF["DF"].apply(lambda x: darstellung_dict[x])
    xn = range(len(x)) #brauchen wir, weil Matplot keine Strings zur X-Achse verarbeiten kann, nur Zahlen
    y = df_DF["Artikel-Lesewert (Erscheinung) in %"]
    yn = df_DF["SplitId"]
    
    #Zeichenfunktion wird aufgerufen
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.4, article = "total", grid=False,\
                  title_text = "Darstellungsformen und Lesewerte")
    
    return prs



#%% Artikellänge und Lesewert/Artikelanzahl
# TODO: Noch nicht geschrieben
   
   
   

#%%





#%% ANALYSE RESSORTS
   
#%%  Initial Ressort-Analyse
# TODO: def analyse_ressorts(prs, df):


    
#%%  Deckblätter Function Zwischenbericht
    ''' Diese Funktion erstellt die Deckblätter inklusive der Screenshots
    für die Auswertungen (einzelne Ressorts, Lokalteile etc.).  
    
    '''
    
def deckblatt_macher(prs, df, ressort,  platzierung, darstellungsform, 
                     seitentitel, ZTG= "null"):

    
    
   
    
    #neuen Slide aufrufen
    
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
   
    ressort_neu = ressort.replace("/", "_").replace(" ", "_").replace("&", "und")
    
    if ZTG !="null": 
        page_path = "./seiten_schwaebische/Schw_" + ZTG + "_" + ressort_neu + ".jpg"
        
    elif ZTG=="null": 
        page_path = "./seiten_schwaebische/Schw_" + ressort_neu + ".jpg"
       
        
    page_left = Inches(7)
    page_top = Inches(0.9)
    page_width = Inches(2.7)
    try:
        page = slide.shapes.add_picture(page_path, page_left, page_top, width=page_width)
    except FileNotFoundError:
        print(ressort_neu + " ... JPEG Titelseite nicht gefunden")
    
    #Form der Tabelle festlegen
        
    rows=1+len(platzierung)+len(darstellungsform)+len(seitentitel)+1
    if len(seitentitel)>0:
        rows +=1
   

    
    cols = 4
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (5)
    height = Inches(0.6)
    
    total_rows = 0

    
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title
    if ZTG != "null":
        title_placeholder.text = ressort + " - " + df.iloc[0]["ZTG"]
    elif ZTG == "null":
        title_placeholder.text = ressort 
    
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
    #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird... 
    if rows>1:
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(2.6)
        table.columns[1].width = Inches(0.8)
        table.columns[2].width = Inches(0.8)
        table.columns[3].width = Inches(0.8)
        table.cell(0, 0).text = ""
        table.cell(0, 1).text = "LW"
        table.cell(0, 2).text = "BW"    
        table.cell(0, 3).text = "DW"    
        table.last_row=True 
    
        #Datensatz für die jeweilige Tabelle anfertigen
        #Wir haben jetzt das df einer Zeitung mit den Artikeln eines Ressorts: 
        data = df.copy()
    
    
        # for-Schleife zündet nur, wenn len(seitentitel) != 0 
        for h in range(len(seitentitel)):
            #table.cell(h+1, 0).text = seitentitel_dict[seitentitel[h]]
            table.cell(h+1, 0).text = seitentitel[h]
            table.cell(h+1, 1).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Lesewert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(h+1, 2).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Blickwert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(h+1, 3).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1)).replace(".", ",")
        #if h==(len(seitentitel)-1):
            #table.cell(h)
            
            total_rows = len(seitentitel)+2
            
        for i in range(len(platzierung)):
            x=0
            y = len(platzierung)+1
            if len(seitentitel) >0: 
                y = y+len(seitentitel)+1
            if len(seitentitel) >= 1:
                x = len(seitentitel)+1
            table.cell(i+1+x,0).text = platzierung_dict[platzierung[i].strip()]
            table.cell(i+1+x,1).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1+x,2).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1+x,3).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            # TODO: Richtige Zeilennummer ausrechnen, falls keine Darstellungsform abgefragt wird
            total_rows = len(platzierung)+1
            if len(seitentitel) >= 1:
                total_rows = len(platzierung)+1+len(seitentitel)+1
            
        for k in range(len(darstellungsform)):
            y = len(platzierung)+1
            if len(seitentitel) >0: 
                y = y+len(seitentitel)+1
            table.cell(k+y,0).text  = darstellung_dict[darstellungsform[k]]
            table.cell(k+y,1).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Lesewert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(k+y,2).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Blickwert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(k+y,3).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1)).replace(".", ",")
            
            total_rows = k+y+1
        
        table.cell(total_rows, 0).text = "Gesamt"
        table.cell(total_rows, 1).text = str(round(data["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(total_rows, 2).text = str(round(data['Artikel-Blickwert (Erscheinung) in %'].mean(), 1)).replace(".", ",")
        table.cell(total_rows, 3).text = str(round(data["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
    
    
    # Schrift in Campton umwandeln etc.  
    
    # Schrift in erster Zeile (Index) umstellen
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.CENTER
    
    for i in range(rows):
        for j in range (cols):
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
   
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,0).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
                
        
        #  Schrift verkleinern, wenn mehr als sieben Zeilen vorhanden sind
        #if rows > 7: 
        #    font_size = 10 
        #else:
        #    font_size = 12
        #for row in range(rows):
        #    for col in range(cols):
        #        table.cell(row, col).text_frame.paragraphs[0].font.size=Pt(font_size)
                
        
        
        
    return prs



#%% Deckblatt-Macher für Abschlussbericht
def deckblatt_abschluss(prs, df_):
    lw_mean = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
    dw_mean = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    bw_mean = df_["Artikel-Blickwert (Erscheinung) in %"].mean()
    
    # Anlegen des Powerpoint-Dokuments
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_placeholder = slide.shapes.title
    title_placeholder.text= df_.iloc[0]["Ressortbeschreibung"]
    placeholder_object = slide.placeholders[10] #unterer Text, größtenteils vom Coach auszufüllen
    placeholder_body = slide.placeholders[12] # oberer Text, orange Box mit LW
    
    placeholder_body.text = "Lesewert Ø " + '{:1.1f}'.format(lw_mean).replace(".", ",")
    placeholder_object.text = "Blickwert: " + '{:1.1f}'.format(bw_mean).replace(".", ",") + ", Durchlesewert: " + '{:1.1f}'.format(dw_mean).replace(".", ",")

    return prs


#%% Tabelle Darstellungsformen / Platzierungen + Seitenbild
    
#Unter research_object verbirgt sich der Untersuchungsgegenstand
# (Darstellungsformen oder Platzierungen)
    
# Unter sort kann eingestellt werden, ob die Ergebnisse nach Seitennummer oder 
    # nach Lesewert sortiert werden sollen. 

def tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", sort="Lesewert"):
    
    print("Funktion tabelle_ressortauswertungt läuft")
    df_ = df.copy()
    df_ = df_[~df_[research_object].isnull()]
    
        
    if research_object == "Darstellungsformen":
        # alle Zeilen mit nur einem Wert
        df_oneval = df_[df_[research_object].str.len()<=2]
        print(df_oneval["Darstellungsformen"].nunique())
        # alle Zeilen mit zwei Werten, hier erster Wert ausgewähl
        df_firstval = df_[df_[research_object].map(len)>2]
        df_firstval["Darstellungsformen"] = df_firstval["Darstellungsformen"].apply(lambda x: x[:2])
        print(df_firstval["Darstellungsformen"].nunique())
        # alle Zeilen mit zwei Werten, hier zweiter Wert ausgewählt
        df_secondval = df_[df_[research_object].map(len)>2]
        df_secondval["Darstellungsformen"] = df_secondval["Darstellungsformen"].apply(lambda x: x[-2:])
        print(df_secondval["Darstellungsformen"].nunique())
        df_ = df_oneval.append([df_firstval, df_secondval], ignore_index=True)
        element_list = darstellungsform
        element_dict = darstellung_dict
    elif research_object == "Platzierungen": 
        element_list = platzierung
        element_dict = platzierung_dict
    
    df_ = df_[df_[research_object].isin(element_list)]
    
    df_ = df_.groupby(research_object, as_index=False)
    df_ = df_.agg({"Seitennummer":"median", "SplitId":"size", 
                   "Artikel-Lesewert (Erscheinung) in %":"mean", 
                   "Artikel-Blickwert (Erscheinung) in %":"mean", 
                   "Artikel-Durchlesewerte (Erscheinung) in %":"mean"})
    
    
    
    if sort=="Seitennummer":
        df_ = df_.sort_values(by="Seitennummer")
    elif sort == "Lesewert": 
        df_ = df_.sort_values(by="Artikel-Lesewert (Erscheinung) in %")
    
    # Tabelle einzeichnen
    print("Tabellen werden angefertigt")   
    # neues Sheet aufrufen
    slide_layout = prs.slide_layouts[2] #ist im Master Überschrift inkl aller Symbole, ansonsten leer
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    # Form der Tabelle Festlegen
    rows = df_.shape[0]+1
    # Nummer der Spalten abhängig von Gesamt- oder Ressort-Auswertung
    cols = 5
    
    left = Inches(1.05)
    top = Inches(1.55)
    width = Inches (5.8)
    height = Inches(0.6)
    
    # jetzt legen wir die Tabelle an, ein table-Objekt
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    
    # Breite der Spalten festlegen (für modus Gesamt und Ressort)
    
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(0.7)
    table.columns[2].width = Inches(0.7)
    table.columns[3].width = Inches(0.7)
    table.columns[4].width = Inches(1)
    
    
    table.cell(0, 0).text = research_object[:-2]
    table.cell(0, 1).text = "LW"
    table.cell(0, 2).text = "BW"
    table.cell(0, 3).text = "DW"
    table.cell(0, 4).text = "Anzahl"
    
    
    # Überschrift festlegen
    title_placeholder = slide.shapes.title
    title = research_object
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
        
    print("Tabelle wird gefüllt")
    
    # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
        # table.last_row = True
    
    
    # data rows, Tabelle füllen
    for i in range(0, df_.shape[0]):
        print("for-Schleife Tabellenfüllung")
        print(df_.iloc[i][research_object])
        table.cell(i+1, 0).text = element_dict[df_.iloc[i][research_object]]
        table.cell(i+1, 1).text = str(round(df_.iloc[i]["Artikel-Lesewert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 2).text = str(round(df_.iloc[i]["Artikel-Blickwert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 3).text = str(round(df_.iloc[i]["Artikel-Durchlesewerte (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 4).text = str(df_.iloc[i]["SplitId"])
    
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i+1, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,0).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
   
    return prs
    
    


#%% Top-10 Function
    
''' 
Die Top-10 Helper-Function erstellt eine Tabelle mit den zehn besten Artikeln

Bei df_berechnet = False errechnet sich die Funktion die zehn besten Artikel 
anhand des Lesewerts selbst. Ansonsten = True bedeutet, dass das Dataframe
bereits in der richtigen Reihenfolge und mit der richtigen Anzahl an Datensätzen 
eingespielt wird. 

überschrift = "default" bedeutet, dass das Programm sich die Überschrift selbst
aus dem Datensatz erstellt. In Einzelfällen (wie die 10 besten Artikel Gesamt),
kann sie auch händisch eingegeben werden.

Mode = "ressort" bedeutet, dass nur Artikel aus einem Ressort verwertet werden, 
eine zusätzlich Anzeige des Ressorts in der Tabelle ist also unnötig.  
Mode = "gesamt" bedeutet, dass in der Tabelle eine Spalte mit der Bezeichnung des
Ressorts eingebaut wird. 

Screenshots = True bedeutet, dass die besten Top5-Artikel als Screenshots 
eingebaut werden. Die Zahl kann unter number_screenshots verändert werden. 

Bei Zeitung = True wird der Zeitungstitel mit in der ÜS ausgespielt. 
 

'''

def top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           mode="ressort", headline="default", zeitung=True):
    
    # falls Top Ten noch nicht berechnet, hier berechnen
    if df_berechnet == False:
        df = df.sort_values(by="Artikel-Lesewert (Erscheinung) in %", 
                            ascending=False).head(10)
        
    # neues Sheet aufrufen
    slide_layout = prs.slide_layouts[2] #ist im Master Überschrift inkl aller Symbole, ansonsten leer
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    # Form der Tabelle Festlegen
    rows = df.shape[0]+1
    # Nummer der Spalten abhängig von Gesamt- oder Ressort-Auswertung
    if mode == "ressort":
        cols = 7
    if mode == "gesamt": 
        cols = 7
    
    left = Inches(1.05)
    top = Inches(1.55)
    width = Inches (8.3)
    height = Inches(0.6)
    
    # jetzt legen wir die Tabelle an, ein table-Objekt
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    
    # Breite der Spalten festlegen (für modus Gesamt und Ressort)
    if mode == "ressort": 
        table.columns[0].width = Inches(0.4)
        table.columns[1].width = Inches(4.0)
        table.columns[2].width = Inches(1.1)
        table.columns[3].width = Inches(0.7)
        table.columns[4].width = Inches(0.7)
        table.columns[5].width = Inches(0.7)
        table.columns[6].width = Inches(0.7)
    
    if mode == "gesamt": 
        table.columns[0].width = Inches(0.4)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(1.1)
        table.columns[3].width = Inches(1.8)
        table.columns[4].width = Inches(0.5)
        table.columns[5].width = Inches(0.5)
        table.columns[6].width = Inches(0.5)
    #  Index benennen
    
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Artikel"
    table.cell(0, 2).text = "Datum"
    # Vierte Spalte wird bei Gesamt zur Ressortspalte
    if mode == "ressort": 
        table.cell(0,3).text = "Seite"
    if mode == "gesamt": 
        table.cell(0,3).text = "Ressort"
        
    table.cell(0, 4).text = "LW"
    table.cell(0, 5).text = "BW"
    table.cell(0, 6).text = "DW"
    
    # Überschrift festlegen
    
    title_placeholder = slide.shapes.title
    title = ""
    if headline == "default":
        if zeitung == True:
            
            title = "Top 10 - " + df.iloc[0].loc["Ressortbeschreibung"] + " " + df.iloc[0].loc["ZTG"]
        elif zeitung == False: 
            title = "Top 10 - " + df.iloc[0].loc["Ressortbeschreibung"]
    
    else:
        title = headline
    
    
    
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    
        
    
    
    # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
        # table.last_row = True
    
    
    # data rows, Tabelle füllen
    for i in range(0, df.shape[0]):
        
        table.cell(i+1, 0).text = str(i+1)
    
        # Checken wie lang der Überschrift-Text ist, bei Bedarf nach dem letzten vollst Wort abschneiden und ... setzen
        if len(df.iloc[i]["Ueberschrifttext"]) > 45:
            kurze_üs = textwrap.shorten(df.iloc[i]["Ueberschrifttext"], width=46, placeholder="...")
            table.cell(i+1, 1).text = kurze_üs
        elif len(df.iloc[i]["Ueberschrifttext"]) <=1:
            table.cell(i+1,1).text = "ACHTUNG!!! KEINE ÜS! SplitId: " + df.iloc[i]["SplitId"]
        else: 
            table.cell(i+1, 1).text = df.iloc[i]["Ueberschrifttext"][:47]
        # jetzt aus dem Timestamp eine lesbare Datumsangabe machen
        datum = df.iloc[i]["Erscheinungsdatum"].strftime("%d.%m.%Y")
        
        table.cell(i+1, 2).text = datum
        
        # Beschrift bei Gesamt= REssort, bei Ressortauswertung = Seitennummer
        if mode=="ressort": 
            table.cell(i+1, 3).text = str(df.iloc[i]["Seitennummer"])
        if mode=="gesamt": 
            table.cell(i+1, 3).text = df.iloc[i]["Ressortbeschreibung"]
        
        table.cell(i+1, 4).text = str(round(df.iloc[i]["Artikel-Lesewert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 5).text = str(round(df.iloc[i]["Artikel-Blickwert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 6).text = str(round(df.iloc[i]["Artikel-Durchlesewerte (Erscheinung) in %"], 1)).replace(".", ",")
    
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i+1, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
   
    screenshot(prs, df, zeitung=zeitung)
        
    return prs 
    
    
    
    
#%% Screenshot-Function    
    
'''
screenshot() fertigt die Screenshots inklusive der LW-Marken an. 
Sie kann händisch aufgerufen oder direkt über die function Top-10 aktiviert 
werden. 

Rangliste = False bedeutet, dass die oberste Ziffer der Marke offen bleibt. 
Rangliste = True bedeutet, dass die obere Ziffer aus df["Rangliste"] bedient wird, 
also die Platzierung des Artikels unter allen Artikeln des Tages. 
Bei vielen Ausgaben ist das aber nur in den Lokalteilen mit den nichtkumulierten 
Artikeln möglich... 

Number of Screenshots (number) ist auf 5 voreingestellt, kann aber verändert werden. 

ausgabe... hier sind Spezialfälle der jeweiligen Ausgabe hinterlegt. Kann 
man händisch ändern oder der Funktion als Parameter mitgeben. 

'''


def screenshot(prs, df, number=5, rangliste=False, ausgabe = "schwaebische", zeitung=True):
    df = df.sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False).head(number)
    for i in range(df.shape[0]):
        #id_nr = "1011" #für die Schwäbische Zeitung
        
        
        elem = df.iloc[i]
        # normale URL zum Laden der Screenshots
        url = "https://lesewert.azureedge.net/layer/"+id_nr+"/"+ str(elem.AusgabeId)+ "/"+ str(elem.ErscheinungsId) + "/"+\
        str(elem.Seitennummer) + "/" +  str(elem.ArtikelId) +  ".jpg"
        # Ausweich-URL zum Laden der Screenshots
        url2= "https://mein.lesewert.de/Offline/ContentProxy?url=" + id_nr + "/"+ str(elem.AusgabeId) + "/"+ \
        str(elem.ErscheinungsId) + "/"+ str(elem.Seitennummer) + "/" +  str(elem.ArtikelId) +  ".jpg"
        
        response = requests.get(url, stream=True)
        response.raw.decode_content = True
    
        # Try-Block prüft, ob URL antwortet
        # Wenn nicht wird alternative URL geladen
        try:
            im = Image.open(response.raw)
        except OSError:
            print("ACHTUNG FEHLER")
            print("Konnte folgende URL nicht laden:")
            print(url)
                       
            response2 = requests.get(url2, stream=True)
            response2.raw.decode_content = True
            im = Image.open(response2.raw)
    
        #Breite des Bildes verringern, wenn der Text höher als breit ist, z.B. ganze Reportage-Seite
        # nur dann ist garantiert, dass die ÜS im Bild ist. 
        if im.size[0] < im.size[1]:
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische":
                height = int(400*corr)
                im = im.resize((400, height), Image.ANTIALIAS)
            else:
                height = int(550*corr)
                im = im.resize((550, height), Image.ANTIALIAS)
        else: 
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische": 
                height = int(600*corr)
                im = im.resize((600, height), Image.ANTIALIAS)
        #TODO:_ Achtung, an den Fotos rumgefummelt
        
        # Bildbearbeitung mit PILLOW
        # Beschriften der LW-Marke
        marke = Image.open("lw_marke_png.png")
        draw = ImageDraw.Draw(marke)
        
        font = ImageFont.truetype('Campton-Light.otf', size=60) 
        font_bold = ImageFont.truetype('Campton-Bold.otf', size=60) 
        (x1,y1) = (340,200)
        (x2,y2) = (340, 300)
        (x3,y3) = (340, 400)
        (x4,y4) = (340, 490)
        message1 = "!!"
        #message1 = str(elem["Platz"]) + str(".")
        message2 = str(round(elem["Artikel-Lesewert (Erscheinung) in %"],1)).replace(".", ",") + "%"
        message3 = str(round(elem["Artikel-Blickwert (Erscheinung) in %"],1)).replace(".", ",") + "%"
        message4 = str(round(elem["Artikel-Durchlesewerte (Erscheinung) in %"],1)).replace(".", ",") + "%"
        color="rgb(255, 255, 255)"
        draw.text((x1,y1), message1, fill=color, font=font)
        draw.text((x2,y2), message2, fill=color, font=font_bold)
        draw.text((x3,y3), message3, fill=color, font=font)
        draw.text((x4,y4), message4, fill=color, font=font)
   
        # Größe der Marke anpassen
        marke = marke.resize((944, 708), Image.ANTIALIAS)
            
        
        #Zwischenspeichern der Marke
        final1 = Image.new("RGBA", (1000, 377), (255,255,255,255))
        # TODO Feineinstellung, Ausgang (0x,0y)
        final1.paste(im, (0,0))
        if im.size[0]<im.size[1]:        # ist der Artikel lang, wird er kleiner, Marke rutscht nach rechts oben  
            final1.paste(marke, (650,5), marke)
        else: # ansonsten bleibt die Marke unten rechts
            final1.paste(marke, (530,35), marke)

        final1.save("final1.png")  
    
    
        # neuen Slide mit Screenshot des Artikels anlegen
        # prs wurde mit Funktionsaufruf übergeben
        slide_layout = prs.slide_layouts[11]  # Layout mit Bild
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
    
        
        title_placeholder = slide.shapes.title
        
        if zeitung == True: 
            title = 'Top ' + str(i+1) + " - " + elem.Ressortbeschreibung + " " + elem.Ausgabenteil
        elif zeitung == False:
            title = 'Top ' + str(i+1) + " - " + elem.Ressortbeschreibung
        
        title_placeholder.text = title
        title_placeholder.text_frame.paragraphs[0].font.bold = True
        title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
        
        placeholder = slide.placeholders[10]
       
        
        #title_placeholder.text_frame.paragraphs[0].font = font_bold
        picture = placeholder.insert_picture('final1.png')
    
    
    
    return prs



#%% Function Ankündigungen

''' 
Diese Funktion stellt Ankündigungen anderen Nachrichten gegenüber. Sie benötigt
ein vorsortiertes DF, den Rest erledigt sie selbst.
TODO: Hier kann man eine grundlegende Tabellen-Funktion überlegen.  
''' 
def lesewert_ankündigungen(prs, df):
    data = df.copy()
    
    #neue Tabelle anlegen
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Tabelle anlegen
    rows = 4 
    cols = 3
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (8.0)
    height = Inches(0.6)
    font_size = 12
        
    total_rows = 0
        
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title            
    title_placeholder.text = "Ankündigungen - "  + data.iloc[0].Ausgabenteil
    title_placeholder.text_frame.paragraphs[0].font.bold = True
        
        #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
        #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird...    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(3.3)
        
        
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Ankündigungen" 
    table.cell(0, 2).text = "alle anderen Texte" 
        
    #table.last_row=True
    table.cell(1,0).text = "LW"
    table.cell(2,0).text = "BW"
    table.cell(3,0).text = "DW"
    
    # TODO Liste muss für schwäbische gekennzeichnet werden
    liste_ak = ["AK", 'AK, BF','AK, NA', 'AK, HG', 'AK, BN', 'BF, AK', 'BB, AK', 'AK, RK', 'NA, AK', 'AK, BB', 'AK, IV', 'AK, RP' ]
    
    df_ak = data[data["Darstellungsformen"].isin(liste_ak)]
    df_notak = data[~data["Darstellungsformen"].isin(liste_ak)]
       

    
    for i in range(1):
            
                
               
            table.cell(i+1,1).text = str(round(df_ak\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1,2).text = str(round(df_notak\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,1).text = str(round(df_ak\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,2).text = str(round(df_notak\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,1).text = str(round(df_ak\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,2).text = str(round(df_notak\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
           
   # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
    # table.last_row = True


    # data rows, Tabelle füllen
    font_size = 12
    for i in range(0, rows):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        
        #for row in range(rows): 
            #text_frame = table.cell(row,1).text_frame
            #p = text_frame.paragraphs[0]
            #p.alignment = PP_ALIGN.LEFT
                
       
    return prs  
    


#%% Übersicht 3 - Lokales Schwäbische
'''
Helper-Funktion für die Schwäbische Zeitung - Ressort Lokales mit Analysen 
auf Seitentitel-Ebene im Lokalen. 

Die Funktion benötigt ein gereinigtes Dataframe, den Rest erledigt sie selbst. 


''' 
  
    

def analyse_lokales(prs, df, df_raw="noch einzusetzen"):
    
    # Liste der drei Ausgaben
    
    ausgabe_liste = ausgaben_liste
    
    st_liste = seitentitel_lokal
    
    
    df_ = df[df["Ressortbeschreibung"]=="Lokales"]
    
    for ausgabe in ausgabe_liste: 
        
        # für jeden Durchlauf ein Dataframe erstellen
        df_ausgabe = df_[df_["ZTG"]==ausgabe]
        print(ausgabe)
        
         # alle Grundwerte im jedem Schleifendurchlauf auf Null setzen
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel=[]
        
        
        
        
        deckblatt_macher(prs, df_ausgabe, "Lokales", platzierung, 
                         darstellungsform, seitentitel, ZTG= ausgabe)
        
        
        top_10(prs, df_ausgabe, df_berechnet = False, screenshots=True, number_screenshots = 5, 
          mode="ressort", headline="Top 10 Lokales " + df_ausgabe["Ausgabenteil"].unique()[0])
        
        
        # Möglichkeit für einzelne keys oder Ausgaben bestimmte Werte einzustellen
        if ausgabe == "FHA": 
            key = "strange"
        else:
            key = "xlarge"
        
        grafik_lesewert(prs, df_ausgabe, target="seitentitel", minimize=5, label_position=key,
                        ressort_liste = st_liste, order="new", legend="xlarge")
        
        
        lesewert_erscheinung(df_ausgabe, prs, title_text="LW nach Wochentagen - " + df_ausgabe.iloc[0].Ausgabenteil)
        
        
        lesewert_ankündigungen(prs, df_ausgabe)
        
    # Jetzt die Übersicht Lokalsport
    
    df_loksport = df[df["Ressortbeschreibung"]=="Lokalsport"]
    
    for ausgabe in ausgabe_liste: 
        
        df_ausgabe_sport = df_loksport[df_loksport["ZTG"]==ausgabe]
        
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel=[]
        
        deckblatt_macher(prs, df_ausgabe_sport, "Lokalsport", platzierung, 
                         darstellungsform, seitentitel, ZTG= ausgabe)
        
        top_10(prs, df_ausgabe_sport, df_berechnet = False, screenshots=True, number_screenshots = 5, 
            mode="ressort", headline="Top 10 Lokalsport " + df_ausgabe_sport["Ausgabenteil"].unique()[0])
        
        #grafik_lesewert(prs, df_ausgabe_sport, target="seitentitel", minimize=5, label_position="xlarge", order="new")
    
        #lesewert_erscheinung(df_ausgabe_sport, prs, title_text="LW nach Wochentagen - " + df_ausgabe.iloc[0].Zeitung)
        
        
    
    return prs


  
    
#%% 1.Lokalseite finden - Schwäbische
    
''' Diese Funktion sucht bei der Schwäbischen Zeitung die erste Lokalseite
in allen drei Ausgaben heraus. 
Läuft normalerweise direkt im Analysecode im Notebook. 

''' 
# Function um Lokale 1 zu identifizieren



def lokale_eins(df):
    ausgaben = ["BIB", "FHA", "RV"]
    list_datum = data2["Erscheinungsdatum"].unique()
    for ausgabe in ausgaben: 
        
        df_ausgabe = df[df["ZTG"]==ausgabe]
       
        for date in list_datum: 
            
            df_ = df_ausgabe[df_ausgabe["Erscheinungsdatum"]==date]
            df_ = df_[df_["Ressortbeschreibung"]== "Lokales"]
            df_ = df_[(df_["Seitentitel"]!="Veranstaltungen") & (df_["Seitentitel"]!="Kirchen")\
                     & (df_["Seitentitel"]!= "Wir in Kreis und Region") & (df_["Seitentitel"]!="Kultur Lokal")\
                     & (df_["Seitentitel"]!= "Region") & (df_["Seitentitel"]!="Umland")\
                     & (df_["Seitentitel"]!="Oberschwaben & Allgäu")]
            
            df_ = df_.sort_values(by="Seitennummer")
            first_page = df_.iloc[0].Seitennummer
            df_ = df_[df_["Seitennummer"]==first_page]
            for i in df_.index: 
                neuer_st = "Lokale Eins " + ausgabe
                df.set_value(i, "Seitentitel", neuer_st)
        
                
            
            
#lokale_eins(data2)
        



#%% Function Tabelle Infokästen

''' 
Diese Funktion schaut sich die Darstellungsform Infokästen (HG)

'''
def tabelle_infokästen(prs, df):
    df_ = df[df["Darstellungsformen"]=="HG"]
    df_not = df[df["Darstellungsformen"]!="HG"]
    
    # Erstellung Tabelle
     #neue Tabelle anlegen
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Tabelle anlegen
    rows = 4 
    cols = 3
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (8.0)
    height = Inches(0.6)
    font_size = 12
        
    total_rows = 0
        
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title            
    title_placeholder.text = "Lesewert Infokästen"
    title_placeholder.text_frame.paragraphs[0].font.bold = True
        
        #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
        #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird...    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(3.3)
        
        
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Infokästen" 
    table.cell(0, 2).text = "alle anderen Artikel" 
        
    #table.last_row=True
    table.cell(1,0).text = "LW"
    table.cell(2,0).text = "BW"
    table.cell(3,0).text = "DW"
    
        
    for i in range(1):
            
                
               
            table.cell(i+1,1).text = str(round(df_["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1,2).text = str(round(df_not\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,1).text = str(round(df_\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,2).text = str(round(df_not\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,1).text = str(round(df_\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,2).text = str(round(df_not\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
           
   # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
    # table.last_row = True


    # data rows, Tabelle füllen
    font_size = 12
    for i in range(0, rows):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            #p.alignment = PP_ALIGN.LEFT
                
    return prs


    


        
#%% Func gelesene vs veröffentlichte Artikel
# TODO muss noch für die Library angepasst werden, ist jetzt erstmal nur rüberkopiert. 
    

def gelesene_artikel(prs, df_scan, df_nichtkum):
    df_s = df_scan.copy()
    df_nk = df_nichtkum.copy()
    
    #ID in String verwandeln
    df_s["TeilartikelVeroeffentlichungenId"] = df_s["TeilartikelVeroeffentlichungenId"].apply(str)
    
    # Umbennen der ID-Spalte
    df_s.rename(columns={"TeilartikelVeroeffentlichungenId":"ArtikelId"}, inplace=True)
    
    # merge mit dem Datensatz, in dem die doppelten Split-IDs noch vorhanden sind
    df_merge = pd.merge(df_s, df_nk, how="left", on="ArtikelId")
    df_m = df_merge.copy()
    
    df_m = df_m[["ArtikelId", "Erscheinungsdatum", "WellenteilnahmenId", "Ressortbeschreibung", "ZTG", "Ausgabename"]]
    
    # Alle Scans ohne entsprechendes Printgegenüber löschen
    # neue Spalte mit Wochentag anlegen (bei manchen Analysen gefordert)
    df_m = df_m[~df_m["Erscheinungsdatum"].isnull()]
    df_m["day"] = df_m["Erscheinungsdatum"].dt.weekday

    
    # Berechnung durschschnittliche Anzahl Artikel pro Tag
    # Anzahl Erscheinungstage:
    anzahl_et = df_m.Erscheinungsdatum.nunique() 

    # Anzahl Leser pro Tag und Ausgabe bzw. Gesamt
    df_group_leserpT_Gesamt = df_m.groupby(["Erscheinungsdatum", "WellenteilnahmenId"], as_index=False).count()
    df_group_leserpT_Gesamt = df_group_leserpT_Gesamt.groupby("Erscheinungsdatum", as_index=False).count()
    df_group_leserpT_Gesamt.rename(columns={"WellenteilnahmenId":"Leser_Tag_gesamt"}, inplace=True)
    df_group_leserpT_Gesamt = df_group_leserpT_Gesamt[["Erscheinungsdatum", "Leser_Tag_gesamt"]]

    df_group_leserpT_Ausgabe = df_m.groupby(["Erscheinungsdatum", "Ausgabename", "WellenteilnahmenId"], as_index=False).count()
    df_group_leserpT_Ausgabe = df_group_leserpT_Ausgabe.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_group_leserpT_Ausgabe.rename(columns={"WellenteilnahmenId":"Leser_Tag_Ausgabe"}, inplace=True)
    df_group_leserpT_Ausgabe = df_group_leserpT_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Leser_Tag_Ausgabe"]]

    # Merge Leser Ausgabe und Gesamt
    df_group_leser = pd.merge(df_group_leserpT_Ausgabe, df_group_leserpT_Gesamt, how="left", on="Erscheinungsdatum")


    # Anzahl gelesene Artikel pro Tag und Ausgabe bzw. Gesamt
    df_group_artikelpT_Gesamt = df_m.groupby("Erscheinungsdatum", as_index=False).count()
    df_group_artikelpT_Gesamt.rename(columns={"ArtikelId":"Scans_Gesamt"}, inplace=True)
    df_group_artikelpT_Gesamt = df_group_artikelpT_Gesamt[["Erscheinungsdatum", "Scans_Gesamt"]]


    df_group_artikelpT_Ausgabe = df_m.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_group_artikelpT_Ausgabe.rename(columns={"ArtikelId":"Scans_Ausgabe"}, inplace=True)
    df_group_artikelpT_Ausgabe = df_group_artikelpT_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Scans_Ausgabe"]]

    # Zusammenführen der Artikel-Count-Dateien
    df_group_scans = pd.merge(df_group_artikelpT_Ausgabe, df_group_artikelpT_Gesamt, how="left", on="Erscheinungsdatum")

    # Anzahl veröffentlichter Artikel 
    df_Artikel_Gesamt = data2.groupby(["Erscheinungsdatum"], as_index=False).count()
    df_Artikel_Gesamt.rename(columns={"SplitId":"Artikelzahl_Tag_Gesamt"}, inplace=True)
    df_Artikel_Gesamt= df_Artikel_Gesamt[["Erscheinungsdatum", "Artikelzahl_Tag_Gesamt"]]

    df_Artikel_Ausgabe = data2.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_Artikel_Ausgabe.rename(columns={"SplitId":"Artikelzahl_Tag_Ausgabe"}, inplace=True)
    df_Artikel_Ausgabe = df_Artikel_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Artikelzahl_Tag_Ausgabe"]]

    #Final merge
    df_group = pd.merge(df_group_artikel, df_group_leser, how="inner",  on=["Erscheinungsdatum", "Ausgabename"])
    df_group_final = pd.merge(df_group, df_Artikel_Ausgabe, on=["Erscheinungsdatum", "Ausgabename"], how="inner")
    
    # Spalten mit Analyse-Zahlen
    df_group_final["Artikel_Leser_Tag_Gesamt"] = df_group_final["Scans_Gesamt"] / df_group_final["Leser_Tag_gesamt"]
    df_group_final["Artikel_Leser_Tag_Ausgabe"] = df_group_final["Scans_Ausgabe"] / df_group_final["Leser_Tag_Ausgabe"] 
    
    # df fpür Grafik ausspielen
    df_grafik = df_group_final.groupby("Ausgabename", as_index=False).mean()
    
    
    
    # GRAFIK PLOTTEN
    
    # Achsen festlegen
    x= df_grafik["Ausgabename"]
    df_grafik["Label"] = df_grafik["Ausgabename"].apply(lambda x: dict_ausgaben[x])
    labels = df_grafik["Label"]
    xn = range(len(x))
    y = df_grafik["Artikel_Leser_Tag_Ausgabe"]
    yn = df_grafik["Artikelzahl_Tag_Ausgabe"]
    set_font_color ="#8c8f91" 
    
     # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
   
    # Barcharts einzeichnen
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Ø gelesene Artikel/Tag")
            
    
                   
                   
        
    ax1.set_ylabel('Ø Artikel/Tag', color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
    # Werte über die Balken setzen
    for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
        height = p.get_height()
        
        txt = '{:1.1f}'.format(height).replace(".", ",")
        ax1.text(p.get_x()+p.get_width()/2., height + 3, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 30, weight = 1000)
        
       
   
    ax1.plot(xn, yn, alpha=0.2) # Linie wird gefadet
    fill = ax1.fill_between(xn, yn, alpha=1, color="#ffcd92", label="Ø Artikel/Tag") #, label=ax2_y_label
                            # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
    ax1.set_ylim(0,yn.max()+1) # manuelles Setzen der Y-Achse rechts, 
                                #damit die auch bei 0 anfängt. 
   
    
    #ax2.set_ylabel(ax2_y_labeltext, color= set_font_color, \
                   #fontproperties=campton_light, labelpad=15) # labelpad =margin

    # Größe der Achsen-Beschriftung festlegen
        #ax2.yaxis.label.set_size(22)
        
    ax1.yaxis.label.set_size(22)
    
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    #
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    
    
    # Legende einbauen
    
    leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill], markerscale=140)
        
        
        
        
  
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    pos_left = 0.08
    pos_right=0.92
    pos_top=0.85
    pos_bottom = 0.4
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_lesewert_plot_two_axis.png"
    plt.savefig(filename)
    
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    title_text = "Artikel - Erschienen und gelesen"
    lw.picture_sheet(prs, filename, title_text=title_text)            
    
    return prs
    
        



#%%

#%% HELPER-FUCNTIONS

#%% Deckblatt-Generator
    # klappt noch nicht so richtig mit dem Design... 
    
def deckblatt(prs, title_text):
    slide_layout = prs.slide_layouts[15]  # Layout mit zwei Platzhaltern idx 10 + 12
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    title_placeholder = slide.shapes.title
    
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs
    
#%% Func Bilder auf PPTX-Folien

'''
HELPER-FUNCTIONEN



POWERPOINT-ERSTELLUNG

picture_sheet(prs, filename)

Diese Function übernimmt das PRS-Objekt und eine erstellte Bilddatei (z.B. über
Matplotlib) und fertigt damit ein neues Sheet an. 


'''
def picture_sheet(prs, filename, title_text):
    
    slide_layout = prs.slide_layouts[11]  # Layout mit Bild
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    placeholder = slide.placeholders[10]
    title_placeholder = slide.shapes.title
    
    pic_path = filename
    pic_left = Inches(0.4)
    pic_top = Inches(2)
    #pic_height = Inches(8)
    pic_width = Inches(9)
    #pic = slide.shapes.add_picture(pic_path, pic_left, pic_top, width=pic_width)
    picture = placeholder.insert_picture(pic_path)
    
    #Titel festlegen
    title_placeholder = slide.shapes.title
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs
        
     
#%% Zwei Bilder auf einer Folie
    
def double_picture_sheet(prs, filename1, filename2, title_text):
    
    slide_layout = prs.slide_layouts[12]  # Layout mit zwei Platzhaltern idx 10 + 12
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    placeholder1 = slide.placeholders[10]
    placeholder2 = slide.placeholders[12]
    title_placeholder = slide.shapes.title
    
    pic_path1 = filename1
    pic_path2 = filename2
    
    pic_left = Inches(0.4)
    pic_top = Inches(2)
    #pic_height = Inches(8)
    pic_width = Inches(9)
    #pic = slide.shapes.add_picture(pic_path, pic_left, pic_top, width=pic_width)
    picture1 = placeholder1.insert_picture(pic_path1)
    picture2 = placeholder2.insert_picture(pic_path2)
    
    #Titel festlegen
    title_placeholder = slide.shapes.title
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs


#%% Func Zeichnen von LW-Bar-Charts
'''
MATPLOTLIB-ERSTELLUNG

plot_two_axis()

Diese Funktion übernimmt vorgefertigte Werte und erstellt daraus eine 
Grafik mit zwei verschiedenen Y-Achsen (links und rechts), einem Balkendiag
und einer Farbfläche im Hintergrund. 

Dafür muss angegeben werden, wie sich x, y, yn zusammensetzen soll. 
Voreingestellt sind hierzu die x, y, yn etc. Werte, die in der übergeordneten
Funktion errechnet werden müssen. Eigentlich sollten sie dann automatisch 
übernommen werden. TESTEN!!!! 

Außerdem muss es möglich sein, die Position der Grafik im Canvas zu verschieben.
Voreingestellt werden die Standardwerte übergeben, der Nutzer kann dann bei 
Bedarf andere Werte einstellen.  

Mit article="total" rechnet die Funktion die Anzahl aller Artikel im 
Erscheinungsraum zusammen. 
article="mean" zeigt die Zahl der Artikle pro Tag im Durchschnitt an. 

Mit grid=False werden die Gridlinien auf der Y-Achse ausgeschaltet. Können 
durch grid=True wieder eingeschaltet werden. 

Parameter axis zeigt an, ob mit nur einer oder Y-zwei Achsen gearbeitet 
werden soll

Wenn mean_line != 0 gesetzt wird, wird eine Linie eingezeichnet, die den 
Durchschnittswert anzeigt.

mit legend="strange" lassen sich besondere formate bearbeiten

'''


def plot_axis(prs, x=0, labels = 0, xn = 0, y = 0, yn = 0,\
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.33, article = "total", grid=False,\
                  title_text="Bitte Titel eingeben", axis=2, mean_line = 0, legend="normal"):
    
    
    #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
    
    # Werte für die Achsen werden festgelegt
    x = x
    xn = xn
    y = y
    yn = yn
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
    if axis==2: 
        ax2 = ax1.twinx() # Zwillings-Y-Achse anlegen, X-Achse wird geteilt

        # Wenn nötig - Gridlinien
        if grid==True: 
            ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    else:
        if grid==True: 
            ax2.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
     # Barcharts einzeichnen
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Lesewert")
            
    
                   
                   
        
    ax1.set_ylabel('Ø Lesewert in Prozent', color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
    # Werte über die Balken setzen
    for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
        height = p.get_height()
        
        txt = '{:1.1f}'.format(height).replace(".", ",")
        ax1.text(p.get_x()+p.get_width()/2., height + 1, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 30, weight = 1000)
        
       
    if axis==2:
        
     # Korrekte Y2-Beschriftung, Label und Legende
     # Parameter article wird hier verarbeiteta
        if article=="total": 
             ax2_y_labeltext = "Artikel im Messzeitraum"
             ax2_y_label ="Artikelanzahl"
        elif article=="mean":
             ax2_y_labeltext ="Ø Anzahl Artikel pro Tag"
             ax2_y_label = "Ø Artikel/Tag"
    
        ax2.plot(xn, yn, alpha=0.2) # Linie wird gefadet
        fill = ax2.fill_between(xn, yn, alpha=1, color="#ffcd92", label=ax2_y_label)
                            # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
        ax2.set_ylim(0,yn.max()+1) # manuelles Setzen der Y-Achse rechts, 
                                #damit die auch bei 0 anfängt. 
   
    
        ax2.set_ylabel(ax2_y_labeltext, color= set_font_color, \
                   fontproperties=campton_light, labelpad=15) # labelpad =margin

    # Größe der Achsen-Beschriftung festlegen
        ax2.yaxis.label.set_size(22)
        
    ax1.yaxis.label.set_size(22)
    
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    if axis==2: 
        ax2.spines["top"].set_visible(False)
        ax2.spines["top"].set_visible(False)
        ax2.spines["bottom"].set_visible(False)
        ax2.spines["left"].set_visible(False)    
        ax2.spines["right"].set_visible(False)
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    # Bei Bedarf Linie mit dem Durchschnitt einziehen
    if mean_line !=0:
        print(mean_line)
        labeltext = "Ø LW Seitentitel: {:1.1f}".format(float(mean_line)).replace(".", ",")
        linie = ax1.axhline(y=mean_line, xmin=0.01, xmax=0.99, color=set_font_color, label=labeltext)
    
    
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.44
    else:
        legend_height = 1.24
    
    if axis==2: 
        ax2.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    
    # Legende einbauen
        if mean_line != 0:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill, linie], markerscale=140)
        else:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill], markerscale=140)
        
        
        
        
        
    elif axis==1:
        
        if mean_line != 0:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, linie], markerscale=140)
        else:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars], markerscale=140)
    
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_lesewert_plot_two_axis.png"
    plt.savefig(filename)
    
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    picture_sheet(prs, filename, title_text=title_text)            
    
    return prs


#%% Schlagworte nach Begriffen durchsuchen 
    
''' Value=<string> legt fest, ob die Funktion nach Sportarten, Vereinen
oder Unternehmen sucht


sort = "LW" oder "Artikelzahl" zeigt an, wie ddie Top 10 erstellt und wie sie georndet werden soll
''' 
    
def schlagworte_finden(prs, df, value, sort= "LW"):
    print("---- hier kommt Schlagwort finden ----")
    print(df)
    #Erstellen ein neues Dict, um daraus später ein DF zu machen
    result_dict = {}
    #Labels festlegen,die später dem Dataframe übergeben werden
    if value == "sportart":
        labels = ["Sportart", "LW", "BW", "DW", "Artikelanzahl"]
        col = "Themen"
        liste = liste_sportarten
    if value == "vereine":
        labels = ["Verein", "LW", "BW", "DW", "Artikelanzahl"]
        col = "Akteure"
        liste = liste_vereine
    # jetzt durchlaufen wir die Liste mit Analysewörtern
    for elem in liste: 
        # check, ob das Analysewort irgendwo im String vorhanden ist
        # nan werden ausgeblendet
        df_elem = df[df[col].str.contains(elem, na=False)]
        #Dict füllen, Analysewort als Key, die Werte als Values
        result_dict[elem] = round(df_elem["Artikel-Lesewert (Erscheinung) in %"].mean(),1), \
        round(df_elem["Artikel-Blickwert (Erscheinung) in %"].mean(),1),\
        round(df_elem["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1), \
        df_elem.shape[0] # letzter Wert shape[0] = Anzahl gefundener Artikel
    
    # pd-Dataframe erstellen, orient="Index" bedeutet, Keys werden Zeilen, Werte werden Spalten
    sw_df = pd.DataFrame.from_dict(result_dict, orient="index").reset_index()
    
    #Spalten umbenennen
    sw_df.columns=labels
    sw_df.loc[sw_df.Sportart=="American Football", ["Sportart"]] = "Am. Football"
    # Werte sortieren, höchster Lesewert zuerst
    sw_df = sw_df.sort_values(by=sort, ascending=False).head(10)
    print(sw_df)
    x_col = sw_df.columns[0]
    x = sw_df[x_col]
    labels = x
    xn = range(len(x))
    y = sw_df["LW"]
    yn = sw_df["Artikelanzahl"]
    title_string = "Lesewert nach " + x_col
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.33, article = "total", grid=False,\
                  title_text = title_string, axis=2)
    

    
    
    
    return prs

#%% Muss-Kann-Soll
def mks(df_ereignis, df):
    df_ = df_ereignis[(df_ereignis["Key"]== "Priority.High") | 
                                  (df_ereignis["Key"] == "Priority.Low") | 
                                    (df_ereignis["Key"] == "Priority.Medium")]
  
    
    df_["TeilartikelVeroeffentlichungsId"] = df_["TeilartikelVeroeffentlichungsId"].astype(int).apply(str)
    
    
    ausgaben = ["BIB", "FHA", "RV"]
    
    
    df = df[df["Ressortbeschreibung"]=="Lokales"]
    
    for aus in ausgaben: 
        df_aus = df[df["ZTG"]==aus]
        df_merge = pd.merge(df_aus, df_, left_on="ArtikelId", right_on="TeilartikelVeroeffentlichungsId", how="left")
        
        df_muss = df_merge[df_merge["Key"]=="Priority.High"].count()
        df_kann = df_merge[df_merge["Key"]=="Priority.Medium"].count()
        df_soll = df_merge[df_merge["Key"]=="Priority.Low"].count()
        total = df_muss["ArtikelId"]+ df_kann["ArtikelId"] + df_soll["ArtikelId"]
        
        muss = df_muss["ArtikelId"] / total * 100
        kann = df_kann["ArtikelId"] / total * 100
        soll = df_soll["ArtikelId"] / total * 100
        print (muss)
        muss_lw = df_merge[df_merge["Key"]=="Priority.High"]["Artikel-Lesewert (Erscheinung) in %"].mean()
        kann_lw = df_merge[df_merge["Key"]=="Priority.Medium"]["Artikel-Lesewert (Erscheinung) in %"].mean()
        soll_lw = df_merge[df_merge["Key"]=="Priority.Low"]["Artikel-Lesewert (Erscheinung) in %"].mean()
        
        print ("-------------------------  " + aus + "  ----------------")
        print("muss für Lok " + aus +": {}  in Prozent: ".format(df_muss["ArtikelId"]))
        print("in Prozent {:1.1f}".format(muss))
        print("kann für Lok " + aus +": {} in Prozent: ".format(df_kann["ArtikelId"]))
        print("in Prozent {:1.1f}".format(kann))
        print("soll für Lok " + aus +": {} in Prozent: ".format(df_soll["ArtikelId"]))
        print("in Prozent {:1.1f}".format(soll))
       
        print ("LW muss für " + aus + ": {:1.1f}".format(muss_lw))
        print ("LW kann für " + aus + ": {:1.1f}".format(kann_lw))
        print ("LW soll für " + aus + ": {:1.1f}".format(soll_lw))
        
    

    
    
    


#%%
#%% ÄLTERE CODESCHNIPSEL


#%% Übersicht 2 - Mantelressorts             

''' übersicht_ressort erstellt die Deckblätter für die einzelnen Ressort, 
ruft dann Top 10 und fünf Screenshots auf

Problem: Diese Funktion muss theoretisch für jede Ausgabe neu geschrieben
werden..

In df_raw sind noch alle Split-IDs enthalten. Wird benötigt, damit die 
Titelseiten einzeln ausgewertet werden können. 
ODER: df_raw = nichtkumulierte Werte

 
'''



def analyse_mantel(prs, df, df_raw="noch einzusetzen", ressort_liste = [], ausgabe_liste= []): 
    
    # check ob Ressortliste händisch angelegt wurde, ansonsten 
    
   
    # Liste mit allen Ressorts des Dataframes wird angelegt
    ressort_liste = mantel_ressorts
    
    
    # Liste mit den drei Ausgaben
    ausgabe_liste = ausgaben_liste
    
    # Liste mit Ressorts wird durchlaufen
    # Anforderungen der Coaches werden individuell in der for-Schleife festgelegt
    
    for ressort in ressort_liste: 
        
        # alle Grundwerte im jedem Schleifendurchlauf auf Null setzen
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel = df[df["Ressortbeschreibung"]==ressort]
        seitentitel = seitentitel.Seitentitel.unique()
        ressort_gefunden = True
        
        #Festlegung, welche Werte analysiert werden sollen
        # Version Schwäbische Zeitung
#        if ressort == "Titel":
#            platzierung = ["AA"]
#            darstellungsform = ["NA"]
#            seitentitel = ["Unterm Strich", "Leitartikel"]
#        elif ressort == "Wir im Süden":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Seite Drei":
#            platzierung = ["AA", "SK"]
#        elif ressort == "Nachrichten & Hintergrund":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Meinung & Dialog": 
#            platzierung = ["AA"]
#            darstellungsform = ["NA"]
#        elif ressort == "Wirtschaft":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Journal":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Kultur":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Sport":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Ratgeber":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Wochenende":
#            seitentitel = ['Wochenende', 'Menschen', 'Lebensart', 'Unterhaltung',
#                           'Szene am Wochenende', 'Meine Seite']
#        
#            
#        else:
#            print ("Ressort " + ressort + " nicht für Auswertung angefordert")
#            ressort_gefunden = False
    
    
       # weitere Spezialanforderungen an einzelne Ressorts
       # zuerst: check ob ein Ressort aus dem Katalog vorhanden ist
        
        if ressort_gefunden: 
            
            # Datensatz mit Daten nur aus dem betreffenden Ressort anlegen
            df_ = df[df["Ressortbeschreibung"]== ressort]
            #df_r = df_raw[df_raw["Ressortbeschreibung"]==ressort]
            
            # Version Schwäbisch Zeitung
            #Spezialfälle
            
#            if ressort == "Titel":
#                
#                for title in ausgabe_liste:
#                    
#                    #Datensatz für jeweiligen Zeitungstitel
#                    df_changed = df_r[df_r["ZTG"]==title]
#                    deckblatt_macher(prs, df_changed, ressort, platzierung, 
#                                     darstellungsform, seitentitel, ZTG=title)
#                    
#                    top_10(prs, df_changed, df_berechnet = False)
#                    
#            elif ressort == "Wochenende":
#                #ausblenden eines Teilstücks eines bereits gewerteten 
#                #Pro und Contra-Artikels
#                df_we = df_[df_["SplitId"]!="23496"]
#                # für func grafik_lesewert wird Liste mit allen Seitentiteln
#                    # erstellt
#                liste_we = df_we["Seitentitel"].unique()
#                
#                deckblatt_macher(prs, df_we, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_we, df_berechnet = False, zeitung=False)
#                print("Wochenende Graifk Lesewert wird aufgerufen")
#                print (liste_we)
#                grafik_lesewert(prs, df_we, target="seitentitel", minimize=0, ressort_liste = liste_we, 
#                                order="ok", legend="xlarge")
#            elif ressort == "Sport": 
#                deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_, df_berechnet = False, zeitung=False)
#                schlagworte_finden(prs, df_, "sportart", sort = "Artikelanzahl")
#            elif ressort== "Ratgeber":
#                #Artiekl ist mit zwei versch. IDs doppelt vorhanden
#                df_ra=df_[df["SplitId"]!="24505"]
#                deckblatt_macher(prs, df_ra, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_ra, df_berechnet = False, zeitung=False)
#            elif ressort == "Seite Drei":
#                # Artikel ist doppelt vorhanden
#                df_sd=df_[df["SplitId"]!="25837"]
#                deckblatt_macher(prs, df_sd, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_sd, df_berechnet = False, zeitung=False)
#            
#            
#            else:
#                # Achtung: Immer df_ übergeben, damit immer nur die Daten des
#                # aktuellen Ressorts verarbeitet werden
#                deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_, df_berechnet = False, zeitung=False)
            
            deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
                    seitentitel, ZTG="null")
            top_10(prs, df_, df_berechnet = False, zeitung=False)

            



    return prs


#%%   
def analyse_mantel_abschluss(prs, df, liste_ressorts=mantel_ressorts):
    print(liste_ressorts)
    for elem in liste_ressorts:
       
        df_ = df[df["Ressortbeschreibung"]==elem]
        
        #Deckblatt 
        deckblatt_abschluss(prs, df_)
        
        # Zusammenstellung Darstellungsformen für Ressort
        tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", sort="Lesewert")
        
        # Zusammenstellung Platzierungen für Ressort
        tabelle_ressortauswertung(prs, df, research_object="Platzierungen", sort="Lesewert")
        
        # Zusammenstellung Seitentitel
        # unter ressort_liste kann noch eine Liste angegenben werden, um nur dort verzeichnete 
        # Seitentitel/Ressorts etc. auszuwerten. 
        grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="normal",
                    ressort_liste=[], special=False, title_text="", order="ok", legend="normal", sort="Seitennummer", 
                    article="total") 
        
        # Überblick über die Entwicklung Lesewert
        #TODO: Mal rausbekommen, was diese Mittelwert-Linie zu bedeuten hat und wie die errechent werden soll... 
        grafik_entwicklung(prs, df, target="Lesewert", mean_line=0, legend="large", grid=True)
        grafik_entwicklung(prs, df, target="Blickwert", mean_line=0, legend="large", grid=True)
        grafik_entwicklung(prs, df, target="Durchlesewert", mean_line=0, legend="large", grid=True)
        # Seitentitel im Überblick
        # Top 10 und Screenshots
        top_10(prs, df_)
    
        

#%% Übersicht 1 - Func Analyse Gesamt   
    
'''
df_scan ist das Dataframe mit den Messdaten aus der ScanAuswertung. 

df_doublesplitid ist ein optionaler Datensatz, bei dem die multiplen Split_ids
nicht entfernt wurden. Manchmal notwendig. 

''' 

    
def analyse_gesamt(prs, df, df_scan=False, df_doublesplitid=False, 
                   df_nichtkum=False):
    
    
    # Demografie - kümmert sich DD drum
    
    # Lesewert in Zahlen - händisch anlegen
    
    # Chart mit Gesamtwert - händisch anlegen
    
    # Chart mit den drei Lokalteilen - händisch anlegen
    
    # Chart mit Mantelteil - händisch anlegen
    # bzw. für einzelne Titel anlegen
    marken_analyse(df, df_doublesplitid=df_doublesplitid, df_nichtkum=df_nichtkum)
    # Übersicht Sonderseiten
   
    # Anzeige Lesezeit
    # TODO ScanAuswertung table runterladen
    #lesezeit(prs, df_scans)
    
       
     
    
    #df_sonderseite = df[df["Seitentitel"].isin(sonderseiten_liste)]
    
    grafik_lesewert(prs, df, target="ressort", minimize=5, label_position="large", 
                    sort = 'Seitennummer', ressort_liste=mantel_ressorts)
    
    # Chart mit Artikelanzahl und LW nach Wochentage
    ##lesewert_erscheinung(df, prs, title_text="Lesewert nach Wochentagen")
    
    # Chart verwendete Darstellungsform und Anzahl Artikel 
    darstellungsformen(prs, df, minimum = 5)
    
    # Chart mit ressort und lesewert
    ##grafik_lesewert(prs, df, target="ressort", minimize=5, special=True, ressort_liste=ressort_list)
    
    # Chart mit Lesezeit
    ##lesezeit(prs, df_scan)
    
    # Vergleich Infokästen
    tabelle_infokästen(prs, df)
    
    
    

#%%
    
#%% EINZELNE ZEITUNGEN - INITIAL-FUNKTIONEN
# TODO:  Funktion schreiben    
    
#  def neue_westfaelische(): 
    
    

    













        
       
